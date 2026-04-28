from __future__ import annotations

import base64
import html
import hashlib
from io import BytesIO
import json
import re
import zipfile
from xml.etree import ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from fastapi import HTTPException, UploadFile, status
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.opc.package import Part
from pptx.util import Emu
from sqlalchemy import func, select
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.core.db import session_scope
from app.models.base import new_id
from app.models.entities import (
    DesignVersion,
    DraftVersion,
    ExportJob,
    OutlineVersion,
    PageBriefVersion,
    Project,
    ProjectMessage,
    ProjectPage,
    RequirementForm,
    ResearchSession,
)
from app.services.background import dispatcher
from app.services.events import append_event
from app.services.generation import GenerationService
from app.services.research import ResearchService

PAGE_STATUS_VALUES = {"empty", "ready", "running", "confirmed", "stale", "failed"}
PROJECT_STAGE_ORDER = {
    "init": 0,
    "outline": 1,
    "search": 2,
    "draft": 3,
    "design": 4,
    "export": 5,
}
WORKFLOW_CONSTRAINTS = [
    {
        "code": "use_env_test_data",
        "label": "使用 .env 测试环境",
        "detail": "当前项目可以直接使用 .env 里的测试环境数据。",
    },
    {
        "code": "db_reset_allowed",
        "label": "允许清库换 schema",
        "detail": "当前没有有价值历史数据，数据库结构变化不需要兼容老版本。",
    },
    {
        "code": "no_meaningless_fallback",
        "label": "禁止无意义 fallback",
        "detail": "不能返回看似成功、实际错误的兜底结果。",
    },
    {
        "code": "remove_dead_code",
        "label": "移除无用代码",
        "detail": "不保留已经失效的旧逻辑和旧分支。",
    },
    {
        "code": "static_ui_reference",
        "label": "参考 static agent UI",
        "detail": "agent 消息展示可参考 static 下已确认的占位样式。",
    },
    {
        "code": "constraints_survive_compression",
        "label": "约束不能被压缩丢失",
        "detail": "多轮消息压缩后也必须保留这些硬约束。",
    },
]

SVG_IMAGE_CONTENT_TYPE = "image/svg+xml"


class SvgImagePart(Part):
    def __init__(
        self,
        *,
        partname,
        package,
        blob: bytes,
        filename: str,
        width_px: float,
        height_px: float,
    ) -> None:
        super().__init__(partname, SVG_IMAGE_CONTENT_TYPE, package, blob)
        self._filename = filename
        self._width_px = width_px
        self._height_px = height_px
        self._sha1 = hashlib.sha1(blob).hexdigest()

    @property
    def desc(self) -> str:
        return self._filename

    @property
    def ext(self) -> str:
        return "svg"

    @property
    def sha1(self) -> str:
        return self._sha1

    def scale(self, scaled_cx: int | None, scaled_cy: int | None) -> tuple[int, int]:
        native_cx, native_cy = self._native_size

        if scaled_cx and scaled_cy:
            return scaled_cx, scaled_cy
        if scaled_cx and not scaled_cy:
            return scaled_cx, int(round(native_cy * (scaled_cx / native_cx)))
        if scaled_cy and not scaled_cx:
            return int(round(native_cx * (scaled_cy / native_cy))), scaled_cy
        return native_cx, native_cy

    @property
    def _native_size(self) -> tuple[int, int]:
        emu_per_inch = 914400
        # SVG 像素按 CSS 规则处理，1px = 1/96in。
        width = int(emu_per_inch * self._width_px / 96)
        height = int(emu_per_inch * self._height_px / 96)
        return Emu(width), Emu(height)


@dataclass
class AgentRunRecorder:
    service: "PptAgentService"
    project: Project
    stage: str
    scope_type: str
    target_page_id: str | None
    title: str
    origin: str
    message_id: str | None = None
    agent_run_id: str = field(default_factory=new_id)
    router_decision: dict[str, Any] | None = None
    step_results: list[dict[str, Any]] = field(default_factory=list)
    next_recommendations: list[dict[str, Any]] = field(default_factory=list)

    def _commit(self) -> None:
        self.service.session.commit()

    def start(self) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="agent.run.started",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={
                "title": self.title,
                "origin": self.origin,
                "message_id": self.message_id,
            },
        )
        self._commit()

    def set_router_decision(self, decision: dict[str, Any]) -> None:
        self.router_decision = decision
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="router.decision",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload=decision,
        )
        self._commit()

    def step_started(self, step_code: str, step_name: str, reason: str) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="action.step.started",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={"step_code": step_code, "step_name": step_name, "reason": reason},
        )
        self._commit()

    def step_progress(
        self,
        step_code: str,
        step_name: str,
        *,
        progress: dict[str, Any] | None = None,
        result: dict[str, Any] | None = None,
    ) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="action.step.progress",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={
                "step_code": step_code,
                "step_name": step_name,
                "status": "running",
                "progress": progress or {},
                "result": result or {},
            },
        )
        self._commit()

    def step_completed(self, step_code: str, step_name: str, result: dict[str, Any] | None = None) -> None:
        step_result = {
            "step_code": step_code,
            "step_name": step_name,
            "status": "completed",
            "result": result or {},
        }
        self.step_results.append(step_result)
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="action.step.completed",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload=step_result,
        )
        self._commit()

    def step_failed(self, step_code: str, step_name: str, error_message: str) -> None:
        step_result = {
            "step_code": step_code,
            "step_name": step_name,
            "status": "failed",
            "error_message": error_message,
        }
        self.step_results.append(step_result)
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="action.step.failed",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload=step_result,
        )
        self._commit()

    def status_changed(self, payload: dict[str, Any]) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="status.changed",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload=payload,
        )
        self._commit()

    def data_updated(self, payload: dict[str, Any]) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="workspace.data.updated",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload=payload,
        )
        self._commit()

    def set_recommendations(self, recommendations: list[dict[str, Any]]) -> None:
        self.next_recommendations = recommendations
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="recommendations.updated",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={"next_recommendations": recommendations},
        )
        self._commit()

    def emit_message(self, message_id: str) -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="agent.message",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={"message_id": message_id},
        )
        self._commit()

    def complete(self, status_text: str = "completed") -> None:
        append_event(
            self.service.session,
            project_id=self.project.id,
            event_type="agent.run.completed",
            stage=self.stage,
            scope_type=self.scope_type,
            target_page_id=self.target_page_id,
            agent_run_id=self.agent_run_id,
            payload={"status": status_text},
        )
        self._commit()


@dataclass
class ImportedTextBlock:
    text: str
    left_px: float
    top_px: float
    width_px: float
    height_px: float
    font_size_px: float
    fill_color: str = "#374151"
    align: str = "start"
    is_title: bool = False


@dataclass
class ImportedImageBlock:
    left_px: float
    top_px: float
    width_px: float
    height_px: float
    data_uri: str


@dataclass
class ImportedVectorBlock:
    kind: str
    left_px: float
    top_px: float
    width_px: float
    height_px: float
    fill_color: str | None = None
    stroke_color: str | None = None
    stroke_width_px: float = 1.0
    x2_px: float | None = None
    y2_px: float | None = None


@dataclass
class ImportedTableCell:
    left_px: float
    top_px: float
    width_px: float
    height_px: float
    text: str
    fill_color: str
    text_color: str
    font_size_px: float
    align: str
    is_bold: bool = False


@dataclass
class ImportedTableBlock:
    left_px: float
    top_px: float
    width_px: float
    height_px: float
    cells: list[ImportedTableCell]
    stroke_color: str = "#CBD5E1"
    stroke_width_px: float = 1.0


@dataclass
class ImportedSlide:
    slide_no: int
    title: str
    bullets: list[str]
    summary_md: str
    page_role: str
    part_title: str | None
    svg_markup: str
    background_color: str


class PptAgentService:
    def __init__(self, session: Session):
        self.session = session
        self.settings = get_settings()
        self.generator = GenerationService()
        self.research = ResearchService(session)

    def list_projects(self, limit: int = 20) -> list[dict[str, Any]]:
        stmt = select(Project).order_by(Project.updated_at.desc()).limit(limit)
        return [self.serialize_project(item) for item in self.session.scalars(stmt)]

    def create_project(self, title: str | None, request_text: str) -> dict[str, Any]:
        project = Project(
            title=title or self.generator.generate_project_title(request_text),
            request_text=request_text,
            current_stage="init",
            workflow_constraints_json={"items": WORKFLOW_CONSTRAINTS},
        )
        self.session.add(project)
        self.session.flush()
        requirement_form = RequirementForm(
            project_id=project.id,
            status="running",
            fixed_items_json=self._build_fixed_fields(),
            answers_json={},
            ai_questions_json=[],
            page_count_options_json=[],
            suggested_actions_json=[],
        )
        self.session.add(requirement_form)
        self.session.flush()
        self._add_message(
            project_id=project.id,
            stage="init",
            scope_type="project",
            role="user",
            content_md=request_text,
        )
        append_event(
            self.session,
            project_id=project.id,
            event_type="project.created",
            stage="init",
            scope_type="project",
            payload={"title": project.title},
        )
        self.session.commit()
        dispatcher.dispatch(run_bootstrap_job, project.id)
        return self.serialize_project(project)

    def import_project(self, *, file: UploadFile, title: str | None = None) -> dict[str, Any]:
        filename = (file.filename or "imported.pptx").strip() or "imported.pptx"
        if Path(filename).suffix.lower() != ".pptx":
            raise HTTPException(status_code=422, detail="当前仅支持导入 .pptx 文件")

        payload = file.file.read()
        if not payload:
            raise HTTPException(status_code=422, detail="上传文件为空")

        try:
            presentation = self._open_presentation(payload)
        except Exception as exc:  # noqa: BLE001
            raise HTTPException(status_code=422, detail=f"PPTX 解析失败: {exc}") from exc

        imported_slides = self._extract_imported_slides(presentation)
        if not imported_slides:
            raise HTTPException(status_code=422, detail="PPTX 中没有可导入的幻灯片")

        resolved_title = title or self._import_project_title(presentation, filename, imported_slides)
        request_text = f"导入已有 PPT：{filename}"
        default_style = self.generator.list_style_options()[0]["style_id"]

        project = Project(
            title=resolved_title,
            request_text=request_text,
            current_stage="design",
            page_count_target=len(imported_slides),
            style_preset=default_style,
            project_metadata_json={
                "is_imported": True,
                "source_filename": filename,
                "slide_count": len(imported_slides),
                "supports_compat_export": True,
            },
            workflow_constraints_json={"items": WORKFLOW_CONSTRAINTS},
        )
        self.session.add(project)
        self.session.flush()

        requirement_form = RequirementForm(
            project_id=project.id,
            status="imported",
            fixed_items_json=self._build_fixed_fields(),
            answers_json={
                "page_count_target": len(imported_slides),
                "style_preset": default_style,
            },
            ai_questions_json=[],
            page_count_options_json=[],
            suggested_actions_json=[],
            init_search_queries_json=[],
            init_search_results_json=[],
            init_corpus_digest_json={},
        )
        self.session.add(requirement_form)
        self.session.flush()

        source_path = self._persist_imported_source_file(project.id, filename, payload)
        project.source_file_path = str(source_path)
        self._rebuild_pages_from_imported_slides(project, imported_slides)

        outline = OutlineVersion(
            project_id=project.id,
            version_no=1,
            status="ready",
            outline_json=self._build_outline_from_imported_slides(imported_slides),
        )
        self.session.add(outline)

        self._add_message(
            project_id=project.id,
            stage="design",
            scope_type="project",
            role="user",
            content_md=f"导入文件：{filename}",
        )
        self._add_message(
            project_id=project.id,
            stage="design",
            scope_type="project",
            role="assistant",
            content_md=f"已导入 {len(imported_slides)} 页 PPT，并生成基础预览。你现在可以直接修改标题、要点、summary，或基于当前内容重生成 draft / design。",
        )
        append_event(
            self.session,
            project_id=project.id,
            event_type="project.imported",
            stage="design",
            scope_type="project",
            payload={"title": project.title, "page_count": len(imported_slides), "filename": filename},
        )
        self.session.commit()
        return self.serialize_project(project)

    def get_project(self, project_id: str) -> dict[str, Any]:
        return self.serialize_project(self._require_project(project_id))

    def list_messages(self, project_id: str) -> list[dict[str, Any]]:
        self._require_project(project_id)
        stmt = (
            select(ProjectMessage)
            .where(ProjectMessage.project_id == project_id)
            .order_by(ProjectMessage.created_at.asc())
        )
        return [self.serialize_message(item) for item in self.session.scalars(stmt)]

    def create_message(
        self,
        *,
        project_id: str,
        scope_type: str,
        target_page_id: str | None,
        ui_surface: str,
        content_md: str,
        attachments: list[dict[str, Any]],
    ) -> dict[str, Any]:
        project = self._require_project(project_id)
        message = self._add_message(
            project_id=project_id,
            stage=project.current_stage,
            scope_type=scope_type,
            target_page_id=target_page_id,
            role="user",
            content_md=content_md,
            structured_payload_json={
                "ui_surface": ui_surface,
                "attachments": attachments,
            },
        )
        self.session.commit()
        dispatcher.dispatch(run_message_job, message.id)
        return self.serialize_message(message)

    def get_requirement_form(self, project_id: str) -> dict[str, Any]:
        project = self._require_project(project_id)
        if not project.requirement_form:
            raise HTTPException(status_code=404, detail="需求单尚未生成")
        return self.serialize_requirement_form(project.requirement_form)

    def submit_requirement_answers(self, project_id: str, answers: list[dict[str, Any]]) -> dict[str, Any]:
        project = self._require_project(project_id)
        requirement_form = self._require_requirement_form(project)
        merged = dict(requirement_form.answers_json or {})
        for item in answers:
            merged[item["question_code"]] = item["value"]
        requirement_form.answers_json = merged
        project.page_count_target = self._coerce_page_count(merged.get("page_count_target"), project.page_count_target)
        if merged.get("style_preset"):
            project.style_preset = str(merged["style_preset"])
        append_event(
            self.session,
            project_id=project.id,
            event_type="requirements.answers_updated",
            stage="init",
            scope_type="project",
            payload={"answers": merged},
        )
        self.session.commit()
        return self.serialize_requirement_form(requirement_form)

    def patch_requirement_answer(self, project_id: str, question_code: str, value: Any) -> dict[str, Any]:
        return self.submit_requirement_answers(
            project_id,
            [{"question_code": question_code, "value": value}],
        )

    def retry_requirement_source(self, project_id: str, source_id: str) -> dict[str, Any]:
        project = self._require_project(project_id)
        requirement_form = self._require_requirement_form(project)
        search_results = self.research.refresh_search_result_cards(requirement_form.init_search_results_json or [])
        source = next((item for item in search_results if item.get("id") == source_id), None)
        if source is None:
            raise HTTPException(status_code=404, detail="指定资料不存在")

        init_collection = self.research.get_or_create_init_collection(project)
        refreshed_source = self.research.retry_search_result_card(
            collection=init_collection,
            search_result=source,
        )
        self.session.refresh(requirement_form)
        latest_search_results = self.research.refresh_search_result_cards(requirement_form.init_search_results_json or [])
        requirement_form.init_search_results_json = [
            refreshed_source if item.get("id") == source_id else item
            for item in latest_search_results
        ]
        requirement_form.init_search_results_json = self.research.refresh_search_result_cards(
            requirement_form.init_search_results_json
        )
        requirement_form.init_corpus_digest_json = self.research.build_collection_digest(init_collection.id)
        append_event(
            self.session,
            project_id=project.id,
            event_type="workspace.data.updated",
            stage="init",
            scope_type="project",
            payload={
                "entity": "requirement_form",
                "update_kind": "init_source_retry",
                "source_id": source_id,
                "read_status": refreshed_source.get("read_status"),
                "vector_status": refreshed_source.get("vector_status"),
                "document_count": requirement_form.init_corpus_digest_json.get("document_count", 0),
                "chunk_count": requirement_form.init_corpus_digest_json.get("chunk_count", 0),
            },
        )
        self.session.commit()
        return self.serialize_requirement_form(requirement_form)

    def retry_page_search_result(self, project_id: str, page_id: str, source_id: str) -> dict[str, Any]:
        project = self._require_project(project_id)
        page = self._require_page(project_id, page_id)
        search_results = self.research.refresh_search_result_cards(page.page_search_results_json or [])
        source = next((item for item in search_results if item.get("id") == source_id), None)
        if source is None:
            raise HTTPException(status_code=404, detail="指定资料不存在")

        collection = self.research.get_or_create_page_collection(project, page)
        refreshed_source = self.research.retry_search_result_card(
            collection=collection,
            search_result=source,
        )
        self.session.refresh(page)
        latest_search_results = self.research.refresh_search_result_cards(page.page_search_results_json or [])
        page.page_search_results_json = [
            refreshed_source if item.get("id") == source_id else item
            for item in latest_search_results
        ]
        page.page_search_results_json = self.research.refresh_search_result_cards(page.page_search_results_json)
        page.page_corpus_digest_json = self.research.build_collection_digest(collection.id)
        page.search_status = "ready" if page.page_corpus_digest_json.get("document_count") else "failed"
        if refreshed_source.get("read_status") == "ready":
            page.summary_status = "stale" if page.page_summary_md else "empty"
            page.draft_status = "stale" if page.current_draft_version_id else "empty"
            page.design_status = "stale" if page.current_design_version_id else "empty"
        if page.current_research_session_id:
            research_session = self.session.get(ResearchSession, page.current_research_session_id)
            if research_session and research_session.page_id == page.id:
                research_session.candidate_sources_json = page.page_search_results_json
                research_session.status = "completed" if page.page_corpus_digest_json.get("document_count") else "failed"
        self._update_artifact_staleness(page)
        append_event(
            self.session,
            project_id=project.id,
            event_type="workspace.data.updated",
            stage=project.current_stage if project.current_stage != "outline" else "search",
            scope_type="page",
            target_page_id=page.id,
            payload={
                "entity": "page",
                "page_id": page.id,
                "update_kind": "search_result_retry",
                "source_id": source_id,
                "read_status": refreshed_source.get("read_status"),
                "vector_status": refreshed_source.get("vector_status"),
                "document_count": page.page_corpus_digest_json.get("document_count", 0),
                "chunk_count": page.page_corpus_digest_json.get("chunk_count", 0),
            },
        )
        self.session.commit()
        return self.serialize_page(page, include_versions=True)

    def create_requirement_question(self, project_id: str, payload: dict[str, Any]) -> dict[str, Any]:
        requirement_form = self._require_requirement_form(self._require_project(project_id))
        questions = [item for item in requirement_form.ai_questions_json if item.get("question_code") != payload["question_code"]]
        questions.append(payload)
        requirement_form.ai_questions_json = questions
        self.session.commit()
        return self.serialize_requirement_form(requirement_form)

    def update_requirement_question(self, project_id: str, question_code: str, payload: dict[str, Any]) -> dict[str, Any]:
        requirement_form = self._require_requirement_form(self._require_project(project_id))
        questions: list[dict[str, Any]] = []
        for item in requirement_form.ai_questions_json:
            if item.get("question_code") != question_code:
                questions.append(item)
                continue
            next_item = dict(item)
            for key, value in payload.items():
                if value is not None:
                    next_item[key] = value
            questions.append(next_item)
        requirement_form.ai_questions_json = questions
        self.session.commit()
        return self.serialize_requirement_form(requirement_form)

    def delete_requirement_question(self, project_id: str, question_code: str) -> dict[str, Any]:
        requirement_form = self._require_requirement_form(self._require_project(project_id))
        requirement_form.ai_questions_json = [
            item for item in requirement_form.ai_questions_json if item.get("question_code") != question_code
        ]
        answers = dict(requirement_form.answers_json or {})
        answers.pop(question_code, None)
        requirement_form.answers_json = answers
        self.session.commit()
        return self.serialize_requirement_form(requirement_form)

    def confirm_requirements(self, project_id: str, note_md: str | None = None) -> dict[str, Any]:
        project = self._require_project(project_id)
        requirement_form = self._require_requirement_form(project)
        self._validate_requirement_form(project, requirement_form)
        if note_md:
            requirement_form.latest_instruction = note_md
        project.current_stage = "outline"
        append_event(
            self.session,
            project_id=project.id,
            event_type="outline.queued",
            stage="outline",
            scope_type="project",
            payload={"project_id": project.id},
        )
        self.session.commit()
        dispatcher.dispatch(run_outline_job, project.id)
        return self.serialize_project(project)

    def upload_background(self, project_id: str, file: UploadFile) -> dict[str, Any]:
        project = self._require_project(project_id)
        suffix = Path(file.filename or "background.bin").suffix or ".bin"
        target = self.settings.background_path / f"{project_id}{suffix}"
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(file.file.read())
        project.background_asset_path = str(target)
        requirement_form = self._require_requirement_form(project)
        answers = dict(requirement_form.answers_json or {})
        answers["background_asset"] = str(target)
        requirement_form.answers_json = answers
        self.session.commit()
        return {"project_id": project.id, "background_asset_path": str(target)}

    def get_outline(self, project_id: str) -> dict[str, Any]:
        outline = self._get_current_outline(project_id)
        if not outline:
            raise HTTPException(status_code=404, detail="大纲尚未生成")
        return self.serialize_outline(outline)

    def list_pages(self, project_id: str) -> list[dict[str, Any]]:
        self._require_project(project_id)
        stmt = select(ProjectPage).where(ProjectPage.project_id == project_id).order_by(ProjectPage.sort_order.asc())
        return [self.serialize_page(item) for item in self.session.scalars(stmt)]

    def get_page(self, project_id: str, page_id: str) -> dict[str, Any]:
        return self.serialize_page(self._require_page(project_id, page_id), include_versions=True)

    def patch_storyboard(self, project_id: str, parts_payload: list[dict[str, Any]]) -> dict[str, Any]:
        project = self._require_project(project_id)
        current_outline = self._get_current_outline(project_id)
        if not current_outline:
            raise HTTPException(status_code=404, detail="当前项目大纲不存在")
        pages = list(
            self.session.scalars(
                select(ProjectPage).where(ProjectPage.project_id == project_id).order_by(ProjectPage.sort_order.asc())
            )
        )
        outline_payload = json.loads(json.dumps(current_outline.outline_json or {}))
        ppt_outline = outline_payload.get("ppt_outline") or {}
        content_pages = [page for page in pages if page.page_role == "content"]
        existing_page_by_id = {page.id: page for page in content_pages}
        requested_existing_ids: list[str] = []
        rebuilt_parts: list[dict[str, Any]] = []
        ordered_content_pages: list[ProjectPage] = []

        used_page_ids: set[str] = set()
        next_page_code_no = 1
        for page in pages:
            match = re.fullmatch(r"page-(\d+)", page.page_code or "")
            if match:
                next_page_code_no = max(next_page_code_no, int(match.group(1)) + 1)

        def allocate_page_code() -> str:
            nonlocal next_page_code_no
            code = f"page-{next_page_code_no:02d}"
            next_page_code_no += 1
            return code

        for part_payload in parts_payload:
            part_title = str(part_payload.get("part_title") or "").strip() or "未命名章节"
            rebuilt_pages: list[dict[str, Any]] = []
            for page_payload in part_payload.get("pages", []):
                page_id = page_payload.get("page_id")
                title = str(page_payload.get("title") or "").strip() or "新内容页"
                content_outline = [str(item).strip() for item in page_payload.get("content_outline", []) if str(item).strip()]

                if page_id:
                    if page_id in used_page_ids:
                        raise HTTPException(status_code=422, detail="storyboard 页面重复，无法重排")
                    page = existing_page_by_id.get(page_id)
                    if not page:
                        raise HTTPException(status_code=422, detail="storyboard 页面不存在，无法重排")
                    used_page_ids.add(page_id)
                    requested_existing_ids.append(page_id)
                    brief = self._get_current_brief(page)
                    current_title = brief.title if brief else ""
                    current_content_outline = brief.content_outline_json if brief else []
                    if current_title != title or current_content_outline != content_outline or page.part_title != part_title:
                        self._new_page_brief_version(
                            page=page,
                            title=title,
                            content_outline=content_outline,
                            section_title=part_title,
                        )
                        self._mark_page_structure_changed(page)
                else:
                    page = ProjectPage(
                        project_id=project.id,
                        page_code=allocate_page_code(),
                        page_role="content",
                        part_title=part_title,
                        sort_order=0,
                        outline_status="ready",
                        search_status="empty",
                        summary_status="empty",
                        draft_status="empty",
                        design_status="empty",
                        page_summary_md="",
                        page_summary_citations_json=[],
                        page_search_queries_json=[],
                        page_search_results_json=[],
                        page_corpus_digest_json={},
                        artifact_staleness_json={},
                    )
                    self.session.add(page)
                    self.session.flush()
                    self._new_page_brief_version(
                        page=page,
                        title=title,
                        content_outline=content_outline,
                        section_title=part_title,
                    )
                    self._mark_page_structure_changed(page)

                ordered_content_pages.append(page)
                rebuilt_pages.append(
                    {
                        "title": title,
                        "content": content_outline,
                    }
                )
            rebuilt_parts.append(
                {
                    "part_title": part_title,
                    "pages": rebuilt_pages,
                }
            )

        requested_existing_id_set = set(requested_existing_ids)
        deleted_content_pages = [page for page in content_pages if page.id not in requested_existing_id_set]
        for page in deleted_content_pages:
            self.session.delete(page)

        ppt_outline["parts"] = rebuilt_parts

        prefix_pages: list[ProjectPage] = []
        suffix_pages: list[ProjectPage] = []
        seen_content = False
        for page in pages:
            if page.page_role == "content":
                seen_content = True
                continue
            if not seen_content:
                prefix_pages.append(page)
            else:
                suffix_pages.append(page)

        ordered_pages = prefix_pages + ordered_content_pages + suffix_pages
        for sort_order, page in enumerate(ordered_pages, start=1):
            page.sort_order = sort_order

        current_outline_changed = current_outline.outline_json != outline_payload
        current_order_ids = [page.id for page in content_pages]
        next_order_ids = [page.id for page in ordered_content_pages]
        if not current_outline_changed and current_order_ids == next_order_ids:
            return {
                "items": [self.serialize_page(page) for page in ordered_pages],
                "outline": self.serialize_outline(current_outline),
            }

        next_outline = OutlineVersion(
            project_id=project.id,
            version_no=current_outline.version_no + 1,
            status="ready",
            outline_json=outline_payload,
        )
        self.session.add(next_outline)
        self.session.commit()
        return {
            "items": [self.serialize_page(page) for page in ordered_pages],
            "outline": self.serialize_outline(next_outline),
        }

    def patch_page_outline(self, project_id: str, page_id: str, payload: dict[str, Any]) -> dict[str, Any]:
        page = self._require_page(project_id, page_id)
        brief = self._new_page_brief_version(
            page=page,
            title=payload["title"],
            content_outline=payload["content_outline"],
            section_title=payload.get("section_title"),
        )
        self._mark_page_structure_changed(page)
        self.session.commit()
        return self.serialize_page(page, include_versions=True)

    def patch_page_summary(self, project_id: str, page_id: str, summary_md: str) -> dict[str, Any]:
        page = self._require_page(project_id, page_id)
        page.page_summary_md = summary_md
        page.summary_status = "ready"
        page.draft_status = "stale" if page.current_draft_version_id else "empty"
        page.design_status = "stale" if page.current_design_version_id else "empty"
        self._update_artifact_staleness(page)
        self.session.commit()
        return self.serialize_page(page, include_versions=True)

    def queue_page_action(
        self,
        project_id: str,
        page_id: str,
        action_type: str,
        *,
        replace_existing: bool = True,
    ) -> dict[str, Any]:
        self._require_page(project_id, page_id)
        agent_run_id = new_id()
        dispatcher.dispatch(run_page_action_job, project_id, page_id, action_type, agent_run_id, replace_existing)
        return {"status": "queued", "agent_run_id": agent_run_id}

    def queue_batch_action(self, project_id: str, action_type: str) -> dict[str, Any]:
        self._require_project(project_id)
        agent_run_id = new_id()
        dispatcher.dispatch(run_batch_action_job, project_id, action_type, agent_run_id)
        return {"status": "queued", "agent_run_id": agent_run_id}

    def get_page_draft(self, project_id: str, page_id: str) -> dict[str, Any]:
        page = self._require_page(project_id, page_id)
        draft = self._get_current_draft(page)
        if not draft:
            raise HTTPException(status_code=404, detail="当前页初稿尚未生成")
        return self.serialize_draft(draft)

    def get_page_design(self, project_id: str, page_id: str) -> dict[str, Any]:
        page = self._require_page(project_id, page_id)
        design = self._get_current_design(page)
        if not design:
            raise HTTPException(status_code=404, detail="当前页设计稿尚未生成")
        return self.serialize_design(design)

    def create_export(self, project_id: str, export_format: str) -> dict[str, Any]:
        project = self._require_project(project_id)
        if export_format == "zip":
            export_path = self._build_export_archive(project.id)
        elif export_format == "pptx_compat":
            export_path = self._build_export_pptx_compat(project)
        elif export_format == "pptx":
            export_path = self._build_export_pptx(project)
        else:
            raise HTTPException(status_code=400, detail="当前仅支持 zip、pptx 或 pptx_compat 导出")
        export_job = ExportJob(project_id=project_id, export_format=export_format, status="completed", file_path=str(export_path))
        self.session.add(export_job)
        self.session.flush()
        self.session.commit()
        return self.serialize_export(export_job)

    def get_export(self, project_id: str, export_id: str) -> dict[str, Any]:
        export = self.session.get(ExportJob, export_id)
        if not export or export.project_id != project_id:
            raise HTTPException(status_code=404, detail="导出任务不存在")
        return self.serialize_export(export)

    def get_export_file_path(self, project_id: str, export_id: str) -> str:
        export = self.session.get(ExportJob, export_id)
        if not export or export.project_id != project_id:
            raise HTTPException(status_code=404, detail="导出任务不存在")
        return export.file_path

    def get_export_download_name(self, project_id: str, export_id: str) -> str:
        export = self.session.get(ExportJob, export_id)
        if not export or export.project_id != project_id:
            raise HTTPException(status_code=404, detail="导出任务不存在")
        project = self._require_project(project_id)
        suffix = ".pptx" if export.export_format in {"pptx", "pptx_compat"} else ".zip"
        return f"{self._slugify_filename(project.title) or 'ppt-agent-export'}{suffix}"

    def _build_compat_export_state(self, project: Project) -> dict[str, Any]:
        project_metadata = project.project_metadata_json or {}
        if not project_metadata.get("is_imported"):
            return {
                "supports_compat_export": False,
                "compat_export_mode": "unavailable",
                "compat_export_notice": None,
            }

        source_path = Path(project.source_file_path or "")
        if not source_path.exists():
            return {
                "supports_compat_export": False,
                "compat_export_mode": "unavailable",
                "compat_export_notice": "原始 PPT 文件不存在，兼容导出不可用。",
            }

        pages = list(
            self.session.scalars(
                select(ProjectPage).where(ProjectPage.project_id == project.id).order_by(ProjectPage.sort_order.asc())
            )
        )
        original_slide_count = int(project_metadata.get("slide_count") or 0)
        if original_slide_count and self._supports_compat_export_layout(pages, original_slide_count):
            return {
                "supports_compat_export": True,
                "compat_export_mode": "preserve",
                "compat_export_notice": "兼容导出会尽量保留未改动原页，改动页覆盖到对应原页。",
            }

        return {
            "supports_compat_export": True,
            "compat_export_mode": "fallback",
            "compat_export_notice": "当前页顺序或数量已变化，兼容导出会自动回退到普通导出。",
        }

    def _serialize_page_preview(self, page: ProjectPage) -> dict[str, Any]:
        design = self._get_current_design(page)
        if design and design.design_svg_markup.strip():
            return {
                "preview_surface": "design",
                "preview_svg_markup": design.design_svg_markup,
            }
        draft = self._get_current_draft(page)
        if draft and draft.draft_svg_markup.strip():
            return {
                "preview_surface": "draft",
                "preview_svg_markup": draft.draft_svg_markup,
            }
        return {
            "preview_surface": "document",
            "preview_svg_markup": None,
        }

    def _serialize_project_preview(self, project: Project) -> dict[str, Any]:
        cover_page = self.session.scalars(
            select(ProjectPage)
            .where(ProjectPage.project_id == project.id, ProjectPage.page_role == "cover")
            .order_by(ProjectPage.sort_order.asc())
            .limit(1)
        ).first()
        if not cover_page:
            return {
                "preview_surface": "fallback",
                "preview_svg_markup": None,
            }

        page_preview = self._serialize_page_preview(cover_page)
        if page_preview["preview_surface"] == "design":
            return page_preview
        if page_preview["preview_surface"] == "draft":
            return page_preview
        return {
            "preview_surface": "fallback",
            "preview_svg_markup": None,
        }

    def serialize_project(self, project: Project) -> dict[str, Any]:
        page_count = self.session.scalar(select(func.count(ProjectPage.id)).where(ProjectPage.project_id == project.id)) or 0
        project_metadata = project.project_metadata_json or {}
        compat_export_state = self._build_compat_export_state(project)
        return {
            "project_id": project.id,
            "title": project.title,
            "request_text": project.request_text,
            "current_stage": project.current_stage,
            "page_count_target": project.page_count_target,
            "style_preset": project.style_preset,
            "background_asset_path": project.background_asset_path,
            "is_imported": bool(project_metadata.get("is_imported")),
            **compat_export_state,
            "workflow_constraints": project.workflow_constraints_json.get("items", []),
            "page_count": page_count,
            "created_at": project.created_at.isoformat(),
            "updated_at": project.updated_at.isoformat(),
            **self._serialize_project_preview(project),
        }

    def serialize_message(self, message: ProjectMessage) -> dict[str, Any]:
        return {
            "id": message.id,
            "project_id": message.project_id,
            "stage": message.stage,
            "scope_type": message.scope_type,
            "target_page_id": message.target_page_id,
            "role": message.role,
            "content_md": message.content_md,
            "structured_payload_json": message.structured_payload_json,
            "created_at": message.created_at.isoformat(),
        }

    def serialize_requirement_form(self, form: RequirementForm) -> dict[str, Any]:
        project = self._require_project(form.project_id)
        init_search_results = self.research.refresh_search_result_cards(form.init_search_results_json or [])
        return {
            "requirement_form": {
                "project_id": form.project_id,
                "status": form.status,
                "workflow_constraints": project.workflow_constraints_json.get("items", []),
                "init_search_queries": form.init_search_queries_json,
                "init_search_results": init_search_results,
                "init_corpus_digest": form.init_corpus_digest_json,
                "page_count_options": form.page_count_options_json,
                "fixed_items": self._build_fixed_fields(),
                "ai_questions": form.ai_questions_json,
                "answers": form.answers_json,
                "suggested_actions": form.suggested_actions_json,
            }
        }

    def serialize_outline(self, outline: OutlineVersion) -> dict[str, Any]:
        return {
            "outline_version_id": outline.id,
            "project_id": outline.project_id,
            "version_no": outline.version_no,
            "status": outline.status,
            "outline": outline.outline_json,
            "created_at": outline.created_at.isoformat(),
            "updated_at": outline.updated_at.isoformat(),
        }

    def serialize_page(self, page: ProjectPage, include_versions: bool = False) -> dict[str, Any]:
        brief = self._get_current_brief(page)
        draft = self._get_current_draft(page)
        design = self._get_current_design(page)
        page_search_results = self.research.refresh_search_result_cards(page.page_search_results_json or [])
        payload = {
            "page_id": page.id,
            "project_id": page.project_id,
            "page_code": page.page_code,
            "page_role": page.page_role,
            "part_title": page.part_title,
            "sort_order": page.sort_order,
            "source_slide_no": page.source_slide_no,
            "title": brief.title if brief else "",
            "content_outline": brief.content_outline_json if brief else [],
            "outline_status": page.outline_status,
            "search_status": page.search_status,
            "summary_status": page.summary_status,
            "draft_status": page.draft_status,
            "design_status": page.design_status,
            "page_search_queries": page.page_search_queries_json,
            "page_search_results": page_search_results,
            "page_corpus_digest": page.page_corpus_digest_json,
            "page_summary_md": page.page_summary_md,
            "page_summary_citations": page.page_summary_citations_json,
            "current_artifact_staleness": page.artifact_staleness_json,
            "current_brief_version_id": page.current_brief_version_id,
            "current_draft_version_id": page.current_draft_version_id,
            "current_design_version_id": page.current_design_version_id,
            "draft_preview_svg_markup": draft.draft_svg_markup if draft and draft.draft_svg_markup.strip() else None,
            "design_preview_svg_markup": design.design_svg_markup if design and design.design_svg_markup.strip() else None,
            "created_at": page.created_at.isoformat(),
            "updated_at": page.updated_at.isoformat(),
            **self._serialize_page_preview(page),
        }
        if include_versions:
            payload["draft"] = self.serialize_draft(draft) if draft else None
            payload["design"] = self.serialize_design(design) if design else None
        return payload

    def serialize_draft(self, draft: DraftVersion) -> dict[str, Any]:
        return {
            "draft_version_id": draft.id,
            "project_id": draft.project_id,
            "page_id": draft.page_id,
            "version_no": draft.version_no,
            "status": draft.status,
            "page_brief_version_id": draft.page_brief_version_id,
            "research_session_id": draft.research_session_id,
            "draft_svg_markup": draft.draft_svg_markup,
            "created_at": draft.created_at.isoformat(),
            "updated_at": draft.updated_at.isoformat(),
        }

    def serialize_design(self, design: DesignVersion) -> dict[str, Any]:
        return {
            "design_version_id": design.id,
            "project_id": design.project_id,
            "page_id": design.page_id,
            "version_no": design.version_no,
            "status": design.status,
            "draft_version_id": design.draft_version_id,
            "style_pack_id": design.style_pack_id,
            "background_asset_path": design.background_asset_path,
            "design_svg_markup": design.design_svg_markup,
            "style_pack": self.generator.get_style_pack(design.style_pack_id),
            "created_at": design.created_at.isoformat(),
            "updated_at": design.updated_at.isoformat(),
        }

    def serialize_export(self, export_job: ExportJob) -> dict[str, Any]:
        return {
            "export_id": export_job.id,
            "project_id": export_job.project_id,
            "export_format": export_job.export_format,
            "status": export_job.status,
            "file_path": export_job.file_path,
            "created_at": export_job.created_at.isoformat(),
            "updated_at": export_job.updated_at.isoformat(),
        }

    def _require_project(self, project_id: str) -> Project:
        project = self.session.get(Project, project_id)
        if not project:
            raise HTTPException(status_code=404, detail="项目不存在")
        return project

    def _require_requirement_form(self, project: Project) -> RequirementForm:
        if not project.requirement_form:
            raise HTTPException(status_code=404, detail="需求单不存在")
        return project.requirement_form

    def _require_page(self, project_id: str, page_id: str) -> ProjectPage:
        page = self.session.get(ProjectPage, page_id)
        if not page or page.project_id != project_id:
            raise HTTPException(status_code=404, detail="页面不存在")
        return page

    def _add_message(
        self,
        *,
        project_id: str,
        stage: str,
        scope_type: str,
        role: str,
        content_md: str,
        target_page_id: str | None = None,
        structured_payload_json: dict[str, Any] | None = None,
    ) -> ProjectMessage:
        message = ProjectMessage(
            project_id=project_id,
            stage=stage,
            scope_type=scope_type,
            target_page_id=target_page_id,
            role=role,
            content_md=content_md,
            structured_payload_json=structured_payload_json or {},
        )
        self.session.add(message)
        self.session.flush()
        return message

    def _build_fixed_fields(self) -> dict[str, Any]:
        return {
            "page_count": {
                "question_code": "page_count_target",
                "allow_custom": True,
            },
            "style_preset": {
                "question_code": "style_preset",
                "options": self.generator.list_style_options(),
                "allow_custom": True,
            },
            "background_asset": {
                "question_code": "background_asset",
                "allow_upload": True,
                "required": False,
            },
        }

    def _coerce_page_count(self, raw_value: Any, fallback: int | None) -> int | None:
        if raw_value is None:
            return fallback
        if isinstance(raw_value, int):
            return raw_value
        if isinstance(raw_value, str):
            stripped = raw_value.strip()
            if stripped.isdigit():
                return int(stripped)
        return fallback

    def _validate_requirement_form(self, project: Project, form: RequirementForm) -> None:
        answers = form.answers_json or {}
        page_count_target = self._coerce_page_count(answers.get("page_count_target"), project.page_count_target)
        if page_count_target is None:
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="页数目标未填写")
        if not answers.get("style_preset") and not project.style_preset:
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="风格预设未选择")
        missing = [
            item["question_code"]
            for item in form.ai_questions_json
            if not str(answers.get(item["question_code"], "")).strip()
        ]
        if missing:
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=f"补充问题未完成: {', '.join(missing)}")
        if not form.init_corpus_digest_json.get("document_count"):
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="初始化资料池尚未建立")

    def _get_current_outline(self, project_id: str) -> OutlineVersion | None:
        stmt = (
            select(OutlineVersion)
            .where(OutlineVersion.project_id == project_id)
            .order_by(OutlineVersion.version_no.desc(), OutlineVersion.created_at.desc())
            .limit(1)
        )
        return self.session.scalars(stmt).first()

    def _get_current_brief(self, page: ProjectPage) -> PageBriefVersion | None:
        if not page.current_brief_version_id:
            return None
        return self.session.get(PageBriefVersion, page.current_brief_version_id)

    def _get_current_draft(self, page: ProjectPage) -> DraftVersion | None:
        if not page.current_draft_version_id:
            return None
        return self.session.get(DraftVersion, page.current_draft_version_id)

    def _get_current_design(self, page: ProjectPage) -> DesignVersion | None:
        if not page.current_design_version_id:
            return None
        return self.session.get(DesignVersion, page.current_design_version_id)

    def _set_project_stage_at_least(self, project: Project, stage: str) -> None:
        if PROJECT_STAGE_ORDER.get(stage, 0) > PROJECT_STAGE_ORDER.get(project.current_stage, 0):
            project.current_stage = stage

    def _open_presentation(self, payload: bytes):
        from pptx import Presentation

        return Presentation(BytesIO(payload))

    def _persist_imported_source_file(self, project_id: str, filename: str, payload: bytes) -> Path:
        target_dir = self.settings.upload_path / project_id
        target_dir.mkdir(parents=True, exist_ok=True)
        target_path = target_dir / f"source{Path(filename).suffix.lower() or '.pptx'}"
        target_path.write_bytes(payload)
        return target_path

    def _import_project_title(self, presentation: Any, filename: str, slides: list[ImportedSlide]) -> str:
        title = ""
        core_properties = getattr(presentation, "core_properties", None)
        if core_properties is not None and getattr(core_properties, "title", None):
            title = str(core_properties.title).strip()
        if not title:
            filename_title = Path(filename).stem.strip().replace("_", " ")
            title = filename_title
        if slides:
            slide_title = slides[0].title.strip()
            if title and not re.fullmatch(r"第\s*\d+\s*页", title):
                return title
            if slide_title and not re.fullmatch(r"第\s*\d+\s*页", slide_title):
                title = slide_title
        return title or "导入的 PPT"

    def _extract_imported_slides(self, presentation: Any) -> list[ImportedSlide]:
        slide_width_px = self._emu_to_px(presentation.slide_width)
        slide_height_px = self._emu_to_px(presentation.slide_height)
        total_slides = len(presentation.slides)
        imported: list[ImportedSlide] = []

        for slide_no, slide in enumerate(presentation.slides, start=1):
            background_color = self._extract_slide_background_color(slide)
            text_blocks = self._extract_slide_text_blocks(slide)
            image_blocks = self._extract_slide_image_blocks(slide)
            vector_blocks = self._extract_slide_vector_blocks(slide)
            table_blocks = self._extract_slide_table_blocks(slide)
            title = self._resolve_imported_slide_title(text_blocks, slide_no)
            bullets = self._resolve_imported_slide_bullets(text_blocks, title)
            role = self._resolve_imported_slide_role(slide_no, total_slides, title)
            imported.append(
                ImportedSlide(
                    slide_no=slide_no,
                    title=title,
                    bullets=bullets,
                    summary_md=self._build_imported_slide_summary(title, bullets, role),
                    page_role=role,
                    part_title="导入内容" if role == "content" else None,
                    svg_markup=self._build_imported_slide_svg(
                        width_px=slide_width_px,
                        height_px=slide_height_px,
                        background_color=background_color,
                        title=title,
                        bullets=bullets,
                        text_blocks=text_blocks,
                        image_blocks=image_blocks,
                        vector_blocks=vector_blocks,
                        table_blocks=table_blocks,
                    ),
                    background_color=background_color,
                )
            )

        if imported and all(item.page_role != "cover" for item in imported):
            imported[0].page_role = "cover"
            imported[0].part_title = None
        return imported

    def _extract_slide_text_blocks(self, slide: Any) -> list[ImportedTextBlock]:
        blocks: list[ImportedTextBlock] = []
        for shape, offset_left_px, offset_top_px in self._iter_slide_shapes(slide):
            if not getattr(shape, "has_text_frame", False):
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
                continue
            text = self._normalize_imported_text(shape.text)
            if not text:
                continue
            blocks.append(
                ImportedTextBlock(
                    text=text,
                    left_px=offset_left_px + self._emu_to_px(getattr(shape, "left", 0)),
                    top_px=offset_top_px + self._emu_to_px(getattr(shape, "top", 0)),
                    width_px=max(self._emu_to_px(getattr(shape, "width", 0)), 80),
                    height_px=max(self._emu_to_px(getattr(shape, "height", 0)), 24),
                    font_size_px=self._extract_shape_font_size_px(shape, 28 if self._shape_is_title(shape) else 16),
                    fill_color=self._extract_shape_font_color(shape, "#111827" if self._shape_is_title(shape) else "#374151"),
                    align=self._extract_shape_text_align(shape),
                    is_title=self._shape_is_title(shape),
                )
            )
        return blocks

    def _extract_slide_image_blocks(self, slide: Any) -> list[ImportedImageBlock]:
        blocks: list[ImportedImageBlock] = []
        for shape, offset_left_px, offset_top_px in self._iter_slide_shapes(slide):
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                continue
            image = getattr(shape, "image", None)
            if image is None or not getattr(image, "blob", None):
                continue
            content_type = getattr(image, "content_type", None) or "image/png"
            blocks.append(
                ImportedImageBlock(
                    left_px=offset_left_px + self._emu_to_px(getattr(shape, "left", 0)),
                    top_px=offset_top_px + self._emu_to_px(getattr(shape, "top", 0)),
                    width_px=max(self._emu_to_px(getattr(shape, "width", 0)), 40),
                    height_px=max(self._emu_to_px(getattr(shape, "height", 0)), 40),
                    data_uri=f"data:{content_type};base64,{base64.b64encode(image.blob).decode('ascii')}",
                )
            )
        return blocks

    def _extract_slide_vector_blocks(self, slide: Any) -> list[ImportedVectorBlock]:
        blocks: list[ImportedVectorBlock] = []
        for shape, offset_left_px, offset_top_px in self._iter_slide_shapes(slide):
            block = self._build_vector_block_from_shape(shape, offset_left_px=offset_left_px, offset_top_px=offset_top_px)
            if block is not None:
                blocks.append(block)
        return blocks

    def _extract_slide_table_blocks(self, slide: Any) -> list[ImportedTableBlock]:
        blocks: list[ImportedTableBlock] = []
        for shape, offset_left_px, offset_top_px in self._iter_slide_shapes(slide):
            if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.TABLE:
                continue
            table_block = self._build_table_block_from_shape(shape, offset_left_px=offset_left_px, offset_top_px=offset_top_px)
            if table_block is not None:
                blocks.append(table_block)
        return blocks

    def _iter_slide_shapes(self, slide: Any):
        yield from self._iter_shape_collection(getattr(slide, "shapes", []), 0.0, 0.0)

    def _iter_shape_collection(self, shapes: Any, base_left_px: float, base_top_px: float):
        ordered_shapes = sorted(shapes, key=lambda item: (getattr(item, "top", 0), getattr(item, "left", 0)))
        for shape in ordered_shapes:
            yield shape, base_left_px, base_top_px
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                next_left_px = base_left_px + self._emu_to_px(getattr(shape, "left", 0))
                next_top_px = base_top_px + self._emu_to_px(getattr(shape, "top", 0))
                yield from self._iter_shape_collection(getattr(shape, "shapes", []), next_left_px, next_top_px)

    def _build_vector_block_from_shape(self, shape: Any, *, offset_left_px: float = 0.0, offset_top_px: float = 0.0) -> ImportedVectorBlock | None:
        shape_type = getattr(shape, "shape_type", None)
        if shape_type in {MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP}:
            return None
        if getattr(shape, "has_text_frame", False) and shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return None

        shape_type_name = str(shape_type).upper()
        left_px = offset_left_px + self._emu_to_px(getattr(shape, "left", 0))
        top_px = offset_top_px + self._emu_to_px(getattr(shape, "top", 0))
        width_px = max(self._emu_to_px(getattr(shape, "width", 0)), 1)
        height_px = max(self._emu_to_px(getattr(shape, "height", 0)), 1)

        if "LINE" in shape_type_name or "CONNECTOR" in shape_type_name:
            stroke_color = self._extract_line_color(shape, "#94A3B8")
            return ImportedVectorBlock(
                kind="line",
                left_px=left_px,
                top_px=top_px,
                width_px=width_px,
                height_px=height_px,
                stroke_color=stroke_color,
                stroke_width_px=self._extract_line_width_px(shape, 1.5),
                x2_px=left_px + width_px,
                y2_px=top_px + height_px,
            )

        if shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return None

        auto_shape_name = str(getattr(shape, "auto_shape_type", "")).upper()
        if "OVAL" in auto_shape_name or "ELLIPSE" in auto_shape_name:
            kind = "ellipse"
        elif "ROUNDED_RECTANGLE" in auto_shape_name:
            kind = "rounded_rect"
        elif "DIAMOND" in auto_shape_name:
            kind = "diamond"
        elif "TRIANGLE" in auto_shape_name:
            kind = "triangle"
        elif "CHEVRON" in auto_shape_name:
            kind = "chevron"
        elif "HEXAGON" in auto_shape_name:
            kind = "hexagon"
        elif "PENTAGON" in auto_shape_name:
            kind = "pentagon"
        elif "PARALLELOGRAM" in auto_shape_name:
            kind = "parallelogram"
        else:
            kind = "rect"
        fill_color = self._extract_fill_color(getattr(shape, "fill", None))
        stroke_color = self._extract_line_color(shape, None)
        if fill_color is None and stroke_color is None:
            return None
        return ImportedVectorBlock(
            kind=kind,
            left_px=left_px,
            top_px=top_px,
            width_px=width_px,
            height_px=height_px,
            fill_color=fill_color,
            stroke_color=stroke_color,
            stroke_width_px=self._extract_line_width_px(shape, 1.0),
        )

    def _build_table_block_from_shape(self, shape: Any, *, offset_left_px: float = 0.0, offset_top_px: float = 0.0) -> ImportedTableBlock | None:
        table = getattr(shape, "table", None)
        if table is None:
            return None

        origin_left_px = offset_left_px + self._emu_to_px(getattr(shape, "left", 0))
        origin_top_px = offset_top_px + self._emu_to_px(getattr(shape, "top", 0))
        cells: list[ImportedTableCell] = []
        current_top_px = origin_top_px

        for row_index, row in enumerate(table.rows):
            row_height_px = max(self._emu_to_px(row.height), 18)
            current_left_px = origin_left_px
            for col_index, column in enumerate(table.columns):
                cell_width_px = max(self._emu_to_px(column.width), 18)
                cell = table.cell(row_index, col_index)
                cell_text = self._normalize_imported_text(cell.text)
                cells.append(
                    ImportedTableCell(
                        left_px=current_left_px,
                        top_px=current_top_px,
                        width_px=cell_width_px,
                        height_px=row_height_px,
                        text=cell_text,
                        fill_color=self._extract_fill_color(getattr(cell, "fill", None)) or "#FFFFFF",
                        text_color=self._extract_cell_text_color(cell, "#111827"),
                        font_size_px=self._extract_cell_font_size_px(cell, 14),
                        align=self._extract_cell_text_align(cell),
                        is_bold=self._extract_cell_is_bold(cell),
                    )
                )
                current_left_px += cell_width_px
            current_top_px += row_height_px

        return ImportedTableBlock(
            left_px=origin_left_px,
            top_px=origin_top_px,
            width_px=max(self._emu_to_px(getattr(shape, "width", 0)), 1),
            height_px=max(self._emu_to_px(getattr(shape, "height", 0)), 1),
            cells=cells,
        )

    def _resolve_imported_slide_title(self, text_blocks: list[ImportedTextBlock], slide_no: int) -> str:
        for block in text_blocks:
            if block.is_title and block.text.strip():
                return block.text.strip()
        for block in text_blocks:
            if block.text.strip():
                return block.text.strip().split("\n", 1)[0][:120]
        return f"第 {slide_no} 页"

    def _resolve_imported_slide_bullets(self, text_blocks: list[ImportedTextBlock], title: str) -> list[str]:
        bullets: list[str] = []
        normalized_title = self._normalize_imported_text(title)
        for block in text_blocks:
            for raw_line in block.text.split("\n"):
                line = self._normalize_imported_text(raw_line)
                if not line or line == normalized_title:
                    continue
                line = re.sub(r"^[\-\u2022\u25E6\u00B7]+\s*", "", line)
                if line and line not in bullets:
                    bullets.append(line[:180])
                if len(bullets) >= 6:
                    return bullets
        return bullets

    def _resolve_imported_slide_role(self, slide_no: int, total_slides: int, title: str) -> str:
        normalized = title.lower()
        if slide_no == 1:
            return "cover"
        if any(keyword in normalized for keyword in ("目录", "contents", "agenda", "table of contents")):
            return "toc"
        if slide_no == total_slides and any(keyword in normalized for keyword in ("谢谢", "thanks", "thank you", "q&a", "qa", "结束")):
            return "end"
        return "content"

    def _build_imported_slide_summary(self, title: str, bullets: list[str], page_role: str) -> str:
        lines = [f"# {title}"]
        if bullets:
            lines.extend([f"- {item}" for item in bullets[:6]])
        elif page_role == "cover":
            lines.append("- 导入封面页")
        elif page_role == "toc":
            lines.append("- 导入目录页")
        elif page_role == "end":
            lines.append("- 导入结束页")
        else:
            lines.append("- 导入内容页")
        return "\n".join(lines)

    def _build_imported_slide_svg(
        self,
        *,
        width_px: float,
        height_px: float,
        background_color: str,
        title: str,
        bullets: list[str],
        text_blocks: list[ImportedTextBlock],
        image_blocks: list[ImportedImageBlock],
        vector_blocks: list[ImportedVectorBlock],
        table_blocks: list[ImportedTableBlock],
    ) -> str:
        svg_parts = [
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{int(round(width_px))}" height="{int(round(height_px))}" viewBox="0 0 {int(round(width_px))} {int(round(height_px))}">',
            f'<rect width="100%" height="100%" fill="{background_color}"/>',
        ]

        for vector in vector_blocks:
            rendered = self._render_imported_vector_block(vector)
            if rendered:
                svg_parts.append(rendered)

        for image in image_blocks:
            svg_parts.append(
                f'<image href="{image.data_uri}" x="{image.left_px:.2f}" y="{image.top_px:.2f}" width="{image.width_px:.2f}" height="{image.height_px:.2f}" preserveAspectRatio="none"/>'
            )

        for table in table_blocks:
            svg_parts.append(self._render_imported_table_block(table))

        if text_blocks:
            for block in text_blocks:
                rendered = self._render_imported_text_block(block)
                if rendered:
                    svg_parts.append(rendered)
        else:
            svg_parts.append(f'<text x="72" y="120" font-size="30" font-weight="700" fill="#111827">{html.escape(title)}</text>')
            for index, item in enumerate(bullets[:8]):
                svg_parts.append(
                    f'<text x="88" y="{180 + index * 32}" font-size="18" fill="#374151">• {html.escape(item)}</text>'
                )

        svg_parts.append("</svg>")
        return "".join(svg_parts)

    def _render_imported_text_block(self, block: ImportedTextBlock) -> str:
        fill = block.fill_color
        font_weight = "700" if block.is_title else "400"
        line_height = max(block.font_size_px * 1.35, 18)
        if block.align == "middle":
            x = block.left_px + block.width_px / 2
        elif block.align == "end":
            x = block.left_px + block.width_px
        else:
            x = block.left_px
        y = block.top_px + block.font_size_px
        lines = [self._normalize_imported_text(item) for item in block.text.split("\n")]
        lines = [item for item in lines if item]
        if not lines:
            return ""

        text_parts = [
            f'<text x="{x:.2f}" y="{y:.2f}" font-size="{block.font_size_px:.2f}" font-weight="{font_weight}" fill="{fill}" text-anchor="{block.align}" xml:space="preserve">'
        ]
        for index, line in enumerate(lines):
            safe_line = html.escape(line)
            if index == 0:
                text_parts.append(f"<tspan>{safe_line}</tspan>")
            else:
                text_parts.append(f'<tspan x="{x:.2f}" dy="{line_height:.2f}">{safe_line}</tspan>')
        text_parts.append("</text>")
        return "".join(text_parts)

    def _render_imported_vector_block(self, block: ImportedVectorBlock) -> str:
        fill = block.fill_color or "none"
        stroke = block.stroke_color or "none"
        stroke_width = max(block.stroke_width_px, 0.75)
        if block.kind == "line":
            return (
                f'<line x1="{block.left_px:.2f}" y1="{block.top_px:.2f}" '
                f'x2="{(block.x2_px if block.x2_px is not None else block.left_px + block.width_px):.2f}" '
                f'y2="{(block.y2_px if block.y2_px is not None else block.top_px + block.height_px):.2f}" '
                f'stroke="{stroke}" stroke-width="{stroke_width:.2f}" stroke-linecap="round"/>'
            )
        if block.kind == "ellipse":
            return (
                f'<ellipse cx="{block.left_px + block.width_px / 2:.2f}" cy="{block.top_px + block.height_px / 2:.2f}" '
                f'rx="{block.width_px / 2:.2f}" ry="{block.height_px / 2:.2f}" '
                f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width:.2f}"/>'
            )
        if block.kind == "rounded_rect":
            radius = min(block.width_px, block.height_px) * 0.12
            return (
                f'<rect x="{block.left_px:.2f}" y="{block.top_px:.2f}" width="{block.width_px:.2f}" height="{block.height_px:.2f}" '
                f'rx="{radius:.2f}" ry="{radius:.2f}" fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width:.2f}"/>'
            )
        if block.kind == "diamond":
            points = [
                (block.left_px + block.width_px / 2, block.top_px),
                (block.left_px + block.width_px, block.top_px + block.height_px / 2),
                (block.left_px + block.width_px / 2, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px / 2),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        if block.kind == "triangle":
            points = [
                (block.left_px + block.width_px / 2, block.top_px),
                (block.left_px + block.width_px, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        if block.kind == "chevron":
            inset = block.width_px * 0.22
            points = [
                (block.left_px, block.top_px),
                (block.left_px + block.width_px - inset, block.top_px),
                (block.left_px + block.width_px, block.top_px + block.height_px / 2),
                (block.left_px + block.width_px - inset, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px),
                (block.left_px + inset, block.top_px + block.height_px / 2),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        if block.kind == "hexagon":
            inset = block.width_px * 0.18
            points = [
                (block.left_px + inset, block.top_px),
                (block.left_px + block.width_px - inset, block.top_px),
                (block.left_px + block.width_px, block.top_px + block.height_px / 2),
                (block.left_px + block.width_px - inset, block.top_px + block.height_px),
                (block.left_px + inset, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px / 2),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        if block.kind == "pentagon":
            points = [
                (block.left_px + block.width_px / 2, block.top_px),
                (block.left_px + block.width_px, block.top_px + block.height_px * 0.38),
                (block.left_px + block.width_px * 0.82, block.top_px + block.height_px),
                (block.left_px + block.width_px * 0.18, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px * 0.38),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        if block.kind == "parallelogram":
            inset = block.width_px * 0.16
            points = [
                (block.left_px + inset, block.top_px),
                (block.left_px + block.width_px, block.top_px),
                (block.left_px + block.width_px - inset, block.top_px + block.height_px),
                (block.left_px, block.top_px + block.height_px),
            ]
            return self._render_polygon(points, fill, stroke, stroke_width)
        return (
            f'<rect x="{block.left_px:.2f}" y="{block.top_px:.2f}" width="{block.width_px:.2f}" height="{block.height_px:.2f}" '
            f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width:.2f}"/>'
        )

    def _render_polygon(self, points: list[tuple[float, float]], fill: str, stroke: str, stroke_width: float) -> str:
        point_text = " ".join(f"{x:.2f},{y:.2f}" for x, y in points)
        return f'<polygon points="{point_text}" fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width:.2f}"/>'

    def _render_imported_table_block(self, table: ImportedTableBlock) -> str:
        parts: list[str] = []
        for cell in table.cells:
            parts.append(
                f'<rect x="{cell.left_px:.2f}" y="{cell.top_px:.2f}" width="{cell.width_px:.2f}" height="{cell.height_px:.2f}" '
                f'fill="{cell.fill_color}" stroke="{table.stroke_color}" stroke-width="{table.stroke_width_px:.2f}"/>'
            )
            if not cell.text:
                continue
            text_block = ImportedTextBlock(
                text=cell.text,
                left_px=cell.left_px + 8,
                top_px=cell.top_px + 6,
                width_px=max(cell.width_px - 16, 8),
                height_px=max(cell.height_px - 12, 8),
                font_size_px=cell.font_size_px,
                fill_color=cell.text_color,
                align="start" if cell.align == "start" else cell.align,
                is_title=cell.is_bold,
            )
            parts.append(self._render_imported_text_block(text_block))
        return "".join(parts)

    def _build_outline_from_imported_slides(self, slides: list[ImportedSlide]) -> dict[str, Any]:
        cover = next((slide for slide in slides if slide.page_role == "cover"), slides[0])
        toc = next((slide for slide in slides if slide.page_role == "toc"), None)
        end = next((slide for slide in reversed(slides) if slide.page_role == "end"), None)
        content_slides = [slide for slide in slides if slide.page_role == "content"]

        return {
            "ppt_outline": {
                "cover": {
                    "title": cover.title,
                    "sub_title": "",
                    "content": cover.bullets,
                },
                "table_of_contents": {
                    "title": toc.title if toc else "目录",
                    "content": toc.bullets if toc else [slide.title for slide in content_slides[:8]],
                },
                "parts": [
                    {
                        "part_title": "导入内容",
                        "pages": [{"title": slide.title, "content": slide.bullets} for slide in content_slides],
                    }
                ],
                "end_page": {
                    "title": end.title if end else "结束",
                    "content": end.bullets if end else [],
                },
            }
        }

    def _rebuild_pages_from_imported_slides(self, project: Project, slides: list[ImportedSlide]) -> None:
        for page in self.session.scalars(select(ProjectPage).where(ProjectPage.project_id == project.id)):
            self.session.delete(page)
        self.session.flush()

        for slide in slides:
            page = ProjectPage(
                project_id=project.id,
                page_code=f"page-{slide.slide_no:02d}",
                page_role=slide.page_role,
                part_title=slide.part_title,
                sort_order=slide.slide_no,
                source_slide_no=slide.slide_no,
                page_metadata_json={
                    "imported": True,
                    "background_color": slide.background_color,
                },
                outline_status="ready",
                search_status="confirmed",
                summary_status="ready",
                draft_status="ready",
                design_status="ready",
                page_summary_md=slide.summary_md,
                page_summary_citations_json=[],
                page_search_queries_json=[],
                page_search_results_json=[],
                page_corpus_digest_json={},
                artifact_staleness_json={},
            )
            self.session.add(page)
            self.session.flush()
            brief = self._new_page_brief_version(
                page=page,
                title=slide.title,
                content_outline=slide.bullets,
                section_title=slide.part_title,
            )
            draft = DraftVersion(
                project_id=project.id,
                page_id=page.id,
                version_no=1,
                status="ready",
                page_brief_version_id=brief.id,
                research_session_id=None,
                draft_svg_markup=slide.svg_markup,
            )
            self.session.add(draft)
            self.session.flush()
            design = DesignVersion(
                project_id=project.id,
                page_id=page.id,
                version_no=1,
                status="ready",
                draft_version_id=draft.id,
                style_pack_id=project.style_preset,
                background_asset_path=None,
                design_svg_markup=slide.svg_markup,
            )
            self.session.add(design)
            self.session.flush()
            page.current_draft_version_id = draft.id
            page.current_design_version_id = design.id
            page.page_metadata_json = {
                **(page.page_metadata_json or {}),
                "imported": True,
                "source_slide_no": slide.slide_no,
                "background_color": slide.background_color,
                "imported_draft_version_id": draft.id,
                "imported_design_version_id": design.id,
            }
            self._update_artifact_staleness(page)

    def _shape_is_title(self, shape: Any) -> bool:
        if not getattr(shape, "is_placeholder", False):
            return False
        placeholder = getattr(shape, "placeholder_format", None)
        placeholder_type = getattr(placeholder, "type", None)
        return placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

    def _extract_slide_background_color(self, slide: Any) -> str:
        background = getattr(slide, "background", None)
        fill = getattr(background, "fill", None)
        return self._extract_fill_color(fill) or "#FFFFFF"

    def _extract_shape_font_size_px(self, shape: Any, fallback_px: float) -> float:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return fallback_px
        for paragraph in text_frame.paragraphs:
            if paragraph.font and paragraph.font.size:
                return max(float(paragraph.font.size.pt) * 96 / 72, 10)
            for run in paragraph.runs:
                if run.font and run.font.size:
                    return max(float(run.font.size.pt) * 96 / 72, 10)
        return fallback_px

    def _extract_cell_font_size_px(self, cell: Any, fallback_px: float) -> float:
        text_frame = getattr(cell, "text_frame", None)
        if text_frame is None:
            return fallback_px
        for paragraph in text_frame.paragraphs:
            if paragraph.font and paragraph.font.size:
                return max(float(paragraph.font.size.pt) * 96 / 72, 10)
            for run in paragraph.runs:
                if run.font and run.font.size:
                    return max(float(run.font.size.pt) * 96 / 72, 10)
        return fallback_px

    def _extract_shape_font_color(self, shape: Any, fallback: str) -> str:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return fallback
        for paragraph in text_frame.paragraphs:
            if paragraph.font:
                color = self._extract_fore_color_value(getattr(paragraph.font, "color", None))
                if color:
                    return color
            for run in paragraph.runs:
                if run.font:
                    color = self._extract_fore_color_value(getattr(run.font, "color", None))
                    if color:
                        return color
        return fallback

    def _extract_cell_text_color(self, cell: Any, fallback: str) -> str:
        text_frame = getattr(cell, "text_frame", None)
        if text_frame is None:
            return fallback
        for paragraph in text_frame.paragraphs:
            if paragraph.font:
                color = self._extract_fore_color_value(getattr(paragraph.font, "color", None))
                if color:
                    return color
            for run in paragraph.runs:
                if run.font:
                    color = self._extract_fore_color_value(getattr(run.font, "color", None))
                    if color:
                        return color
        return fallback

    def _extract_shape_text_align(self, shape: Any) -> str:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            return "start"
        for paragraph in text_frame.paragraphs:
            alignment = str(getattr(paragraph, "alignment", "") or "").upper()
            if "CENTER" in alignment:
                return "middle"
            if "RIGHT" in alignment:
                return "end"
        return "start"

    def _extract_cell_text_align(self, cell: Any) -> str:
        text_frame = getattr(cell, "text_frame", None)
        if text_frame is None:
            return "start"
        for paragraph in text_frame.paragraphs:
            alignment = str(getattr(paragraph, "alignment", "") or "").upper()
            if "CENTER" in alignment:
                return "middle"
            if "RIGHT" in alignment:
                return "end"
        return "start"

    def _extract_cell_is_bold(self, cell: Any) -> bool:
        text_frame = getattr(cell, "text_frame", None)
        if text_frame is None:
            return False
        for paragraph in text_frame.paragraphs:
            if paragraph.font and paragraph.font.bold:
                return True
            for run in paragraph.runs:
                if run.font and run.font.bold:
                    return True
        return False

    def _extract_fill_color(self, fill: Any) -> str | None:
        if fill is None:
            return None
        try:
            fore_color = fill.fore_color
        except Exception:  # noqa: BLE001
            return None
        return self._extract_fore_color_value(fore_color)

    def _extract_line_color(self, shape: Any, fallback: str | None) -> str | None:
        line = getattr(shape, "line", None)
        try:
            color = self._extract_fore_color_value(getattr(line, "color", None))
        except Exception:  # noqa: BLE001
            color = None
        return color or fallback

    def _extract_line_width_px(self, shape: Any, fallback: float) -> float:
        line = getattr(shape, "line", None)
        width = getattr(line, "width", None)
        if width is None:
            return fallback
        return max(self._emu_to_px(width), 0.75)

    def _extract_fore_color_value(self, color: Any) -> str | None:
        if color is None:
            return None
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return f"#{rgb}"
        if getattr(color, "type", None) is None:
            return None
        brightness = getattr(color, "brightness", 0) or 0
        theme_name = str(getattr(color, "theme_color", "") or "").upper()
        return self._theme_color_fallback(theme_name, brightness)

    def _theme_color_fallback(self, theme_name: str, brightness: float = 0.0) -> str | None:
        palette = {
            "ACCENT_1": "#5B8FF9",
            "ACCENT_2": "#61DDAA",
            "ACCENT_3": "#65789B",
            "ACCENT_4": "#F6BD16",
            "ACCENT_5": "#7262FD",
            "ACCENT_6": "#78D3F8",
            "BACKGROUND_1": "#FFFFFF",
            "BACKGROUND_2": "#F3F4F6",
            "TEXT_1": "#111827",
            "TEXT_2": "#374151",
        }
        base = palette.get(theme_name)
        if not base:
            return None
        return self._adjust_hex_brightness(base, brightness)

    def _adjust_hex_brightness(self, hex_color: str, brightness: float) -> str:
        if not hex_color.startswith("#") or len(hex_color) != 7 or brightness == 0:
            return hex_color
        channels = [int(hex_color[index:index + 2], 16) for index in (1, 3, 5)]
        if brightness > 0:
            adjusted = [int(round(channel + (255 - channel) * brightness)) for channel in channels]
        else:
            adjusted = [int(round(channel * (1 + brightness))) for channel in channels]
        adjusted = [max(0, min(255, value)) for value in adjusted]
        return "#{:02X}{:02X}{:02X}".format(*adjusted)

    def _normalize_imported_text(self, value: Any) -> str:
        text = str(value or "")
        text = text.replace("\r", "\n")
        text = re.sub(r"\n{2,}", "\n", text)
        text = "\n".join(line.strip() for line in text.split("\n"))
        return text.strip()

    def _emu_to_px(self, value: Any) -> float:
        return float(value or 0) / 9525.0

    def _new_page_brief_version(
        self,
        *,
        page: ProjectPage,
        title: str,
        content_outline: list[str],
        section_title: str | None,
        status_text: str = "ready",
    ) -> PageBriefVersion:
        current_brief = self._get_current_brief(page)
        next_version = (current_brief.version_no if current_brief else 0) + 1
        brief = PageBriefVersion(
            project_id=page.project_id,
            page_id=page.id,
            version_no=next_version,
            status=status_text,
            section_title=section_title,
            title=title,
            content_outline_json=content_outline,
            content_summary="；".join(content_outline[:2]) or title,
        )
        self.session.add(brief)
        self.session.flush()
        page.current_brief_version_id = brief.id
        page.part_title = section_title
        page.outline_status = "ready"
        return brief

    def _mark_page_structure_changed(self, page: ProjectPage) -> None:
        if page.page_role == "content":
            if page.page_search_results_json:
                page.search_status = "stale"
            else:
                page.search_status = "empty"
            if page.page_summary_md:
                page.summary_status = "stale"
            else:
                page.summary_status = "empty"
        if page.current_draft_version_id:
            page.draft_status = "stale"
        else:
            page.draft_status = "empty"
        if page.current_design_version_id:
            page.design_status = "stale"
        else:
            page.design_status = "empty"
        self._update_artifact_staleness(page)

    def _update_artifact_staleness(self, page: ProjectPage) -> None:
        page.artifact_staleness_json = {
            "search": page.search_status == "stale",
            "summary": page.summary_status == "stale",
            "draft": page.draft_status == "stale",
            "design": page.design_status == "stale",
        }

    def _build_outline_snapshot(self, project_id: str) -> list[dict[str, Any]]:
        snapshot: list[dict[str, Any]] = []
        stmt = select(ProjectPage).where(ProjectPage.project_id == project_id).order_by(ProjectPage.sort_order.asc())
        for page in self.session.scalars(stmt):
            brief = self._get_current_brief(page)
            snapshot.append(
                {
                    "page_id": page.id,
                    "page_code": page.page_code,
                    "page_role": page.page_role,
                    "section_title": page.part_title,
                    "title": brief.title if brief else "",
                    "content_outline": brief.content_outline_json if brief else [],
                }
            )
        return snapshot

    def _build_project_level_status_summary(self, project: Project) -> dict[str, Any]:
        pages = self.list_pages(project.id)
        return {
            "project_stage": project.current_stage,
            "page_count_target": project.page_count_target,
            "style_preset": project.style_preset,
            "pages": [
                {
                    "page_id": item["page_id"],
                    "page_code": item["page_code"],
                    "page_role": item["page_role"],
                    "title": item["title"],
                    "outline_status": item["outline_status"],
                    "search_status": item["search_status"],
                    "summary_status": item["summary_status"],
                    "draft_status": item["draft_status"],
                    "design_status": item["design_status"],
                }
                for item in pages
            ],
        }

    def _build_page_context_for_router(self, page: ProjectPage | None) -> dict[str, Any]:
        if page is None:
            return {}
        brief = self._get_current_brief(page)
        return {
            "page_id": page.id,
            "page_title": brief.title if brief else "",
            "page_bullets": brief.content_outline_json if brief else [],
            "page_section_title": page.part_title,
            "page_outline_status": page.outline_status,
            "page_search_status": page.search_status,
            "page_summary_status": page.summary_status,
            "page_draft_status": page.draft_status,
            "page_design_status": page.design_status,
            "page_search_queries": page.page_search_queries_json,
            "page_corpus_digest": page.page_corpus_digest_json,
            "page_summary_digest": {
                "summary_md": page.page_summary_md,
                "citation_count": len(page.page_summary_citations_json or []),
            },
            "current_artifact_staleness": page.artifact_staleness_json,
            "outline_full_snapshot": self._build_outline_snapshot(page.project_id),
        }

    def _build_export_archive(self, project_id: str) -> Path:
        pages = list(self.session.scalars(select(ProjectPage).where(ProjectPage.project_id == project_id).order_by(ProjectPage.sort_order.asc())))
        export_path = self.settings.export_path / f"{project_id}.zip"
        export_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(export_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            manifest: list[dict[str, Any]] = []
            for page in pages:
                design = self._get_current_design(page)
                if not design:
                    continue
                filename = f"{page.page_code}.svg"
                archive.writestr(filename, design.design_svg_markup)
                manifest.append({"page_id": page.id, "page_code": page.page_code, "file": filename})
            archive.writestr("manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2))
        return export_path

    def _build_export_pptx(self, project: Project) -> Path:
        exportables = self._collect_exportable_designs(project.id)
        from pptx import Presentation

        presentation = Presentation()
        slide_width_emu, slide_height_emu = self._presentation_size_from_svg(exportables[0][1].design_svg_markup)
        presentation.slide_width = slide_width_emu
        presentation.slide_height = slide_height_emu
        blank_layout = presentation.slide_layouts[6]
        svg_part_cache: dict[str, SvgImagePart] = {}

        for page, design in exportables:
            slide = presentation.slides.add_slide(blank_layout)
            left, top, width, height = self._fit_picture_to_slide(
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
                svg_markup=design.design_svg_markup,
            )
            self._add_svg_picture(
                slide=slide,
                svg_markup=design.design_svg_markup,
                filename=f"{page.page_code}.svg",
                left=left,
                top=top,
                width=width,
                height=height,
                cache=svg_part_cache,
            )

        export_path = self.settings.export_path / f"{project.id}.pptx"
        export_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(export_path))
        return export_path

    def _build_export_pptx_compat(self, project: Project) -> Path:
        source_path = Path(project.source_file_path or "")
        if not source_path.exists():
            return self._build_export_pptx(project)

        from pptx import Presentation

        pages = list(
            self.session.scalars(
                select(ProjectPage).where(ProjectPage.project_id == project.id).order_by(ProjectPage.sort_order.asc())
            )
        )
        if not pages:
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail="当前项目没有可导出的页面")

        presentation = Presentation(str(source_path))
        original_slide_count = len(presentation.slides)
        if not self._supports_compat_export_layout(pages, original_slide_count):
            return self._build_export_pptx(project)

        blank_layout = presentation.slide_layouts[6]
        svg_part_cache: dict[str, SvgImagePart] = {}

        while len(presentation.slides) < len(pages):
            presentation.slides.add_slide(blank_layout)

        for index, page in enumerate(pages, start=1):
            slide = presentation.slides[index - 1]
            design = self._get_current_design(page)
            if self._should_keep_original_slide(page):
                continue
            if not design or not design.design_svg_markup.strip():
                raise HTTPException(
                    status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                    detail=f"页面《{self._page_export_title(page)}》缺少可覆盖导出的设计稿",
                )
            left, top, width, height = self._fit_picture_to_slide(
                slide_width_emu=presentation.slide_width,
                slide_height_emu=presentation.slide_height,
                svg_markup=design.design_svg_markup,
            )
            self._add_svg_picture(
                slide=slide,
                svg_markup=design.design_svg_markup,
                filename=f"{page.page_code}.svg",
                left=left,
                top=top,
                width=width,
                height=height,
                cache=svg_part_cache,
            )

        export_path = self.settings.export_path / f"{project.id}.compat.pptx"
        export_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(export_path))
        return export_path

    def _supports_compat_export_layout(self, pages: list[ProjectPage], original_slide_count: int) -> bool:
        if len(pages) < original_slide_count:
            return False
        for expected_index, page in enumerate(pages, start=1):
            if page.source_slide_no is None:
                if expected_index <= original_slide_count:
                    return False
                continue
            if page.source_slide_no != expected_index:
                return False
        return True

    def _should_keep_original_slide(self, page: ProjectPage) -> bool:
        metadata = page.page_metadata_json or {}
        imported_design_version_id = metadata.get("imported_design_version_id")
        if not imported_design_version_id:
            return False
        return page.current_design_version_id == imported_design_version_id

    def _collect_exportable_designs(self, project_id: str) -> list[tuple[ProjectPage, DesignVersion]]:
        pages = list(
            self.session.scalars(
                select(ProjectPage).where(ProjectPage.project_id == project_id).order_by(ProjectPage.sort_order.asc())
            )
        )
        exportables: list[tuple[ProjectPage, DesignVersion]] = []
        missing_pages: list[str] = []
        for page in pages:
            design = self._get_current_design(page)
            if page.design_status != "ready" or not design or not design.design_svg_markup.strip():
                missing_pages.append(f"{page.sort_order}. {self._page_export_title(page)}")
                continue
            exportables.append((page, design))
        if missing_pages:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"以下页面还没有完成设计稿，无法导出 PPTX: {'；'.join(missing_pages)}",
            )
        if not exportables:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail="当前项目没有可导出的设计稿页面",
            )
        return exportables

    def _page_export_title(self, page: ProjectPage) -> str:
        brief = self._get_current_brief(page)
        return brief.title if brief and brief.title.strip() else page.page_code

    def _add_svg_picture(
        self,
        *,
        slide: Any,
        svg_markup: str,
        filename: str,
        left: int,
        top: int,
        width: int,
        height: int,
        cache: dict[str, SvgImagePart],
    ) -> None:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        svg_part = self._get_or_add_svg_part(
            package=slide.part.package,
            svg_markup=svg_markup,
            filename=filename,
            cache=cache,
        )
        relationship_id = slide.part.relate_to(svg_part, RT.IMAGE)
        shape_id = slide.shapes._next_shape_id
        slide.shapes._grpSp.add_pic(shape_id, f"Picture {shape_id - 1}", filename, relationship_id, left, top, width, height)
        slide.shapes._recalculate_extents()

    def _get_or_add_svg_part(
        self,
        *,
        package: Any,
        svg_markup: str,
        filename: str,
        cache: dict[str, SvgImagePart],
    ) -> SvgImagePart:
        svg_blob = svg_markup.encode("utf-8")
        content_hash = hashlib.sha1(svg_blob).hexdigest()
        existing = cache.get(content_hash)
        if existing:
            return existing

        width_px, height_px = self._extract_svg_canvas_size(svg_markup)
        svg_part = SvgImagePart(
            partname=package.next_image_partname("svg"),
            package=package,
            blob=svg_blob,
            filename=filename,
            width_px=width_px,
            height_px=height_px,
        )
        cache[content_hash] = svg_part
        return svg_part

    def _presentation_size_from_svg(self, svg_markup: str) -> tuple[int, int]:
        width_px, height_px = self._extract_svg_canvas_size(svg_markup)
        aspect_ratio = width_px / height_px if width_px and height_px else 16 / 9
        slide_width_inch = 13.333333
        slide_height_inch = slide_width_inch / aspect_ratio if aspect_ratio > 0 else 7.5
        return Emu(int(slide_width_inch * 914400)), Emu(int(slide_height_inch * 914400))

    def _fit_picture_to_slide(self, *, slide_width_emu: int, slide_height_emu: int, svg_markup: str) -> tuple[int, int, int, int]:
        width_px, height_px = self._extract_svg_canvas_size(svg_markup)
        if not width_px or not height_px:
            return 0, 0, slide_width_emu, slide_height_emu
        width_scale = slide_width_emu / width_px
        height_scale = slide_height_emu / height_px
        scale = min(width_scale, height_scale)
        width = int(width_px * scale)
        height = int(height_px * scale)
        left = int((slide_width_emu - width) / 2)
        top = int((slide_height_emu - height) / 2)
        return left, top, width, height

    def _extract_svg_canvas_size(self, svg_markup: str) -> tuple[float, float]:
        try:
            root = ET.fromstring(svg_markup)
        except ET.ParseError:
            return 1600.0, 900.0

        view_box = root.attrib.get("viewBox") or root.attrib.get("viewbox")
        if view_box:
            parts = [part for part in re.split(r"[,\s]+", view_box.strip()) if part]
            if len(parts) == 4:
                try:
                    width = float(parts[2])
                    height = float(parts[3])
                    if width > 0 and height > 0:
                        return width, height
                except ValueError:
                    pass

        width = self._parse_svg_dimension(root.attrib.get("width"))
        height = self._parse_svg_dimension(root.attrib.get("height"))
        if width > 0 and height > 0:
            return width, height
        return 1600.0, 900.0

    def _parse_svg_dimension(self, raw_value: str | None) -> float:
        if not raw_value:
            return 0.0
        match = re.search(r"[-+]?\d*\.?\d+", raw_value)
        if not match:
            return 0.0
        try:
            return float(match.group(0))
        except ValueError:
            return 0.0

    def _slugify_filename(self, raw_value: str) -> str:
        cleaned = re.sub(r"[^\w\u4e00-\u9fff-]+", "-", raw_value.strip(), flags=re.UNICODE)
        cleaned = re.sub(r"-{2,}", "-", cleaned).strip("-")
        return cleaned or "ppt-agent-export"

    def _build_fixed_field_values(self, project: Project, requirement_form: RequirementForm) -> dict[str, Any]:
        answers = requirement_form.answers_json or {}
        return {
            "page_count_target": self._coerce_page_count(answers.get("page_count_target"), project.page_count_target),
            "style_preset": str(answers.get("style_preset") or project.style_preset or ""),
            "background_asset": str(answers.get("background_asset") or project.background_asset_path or ""),
        }

    def _build_fixed_page_summary(self, title: str, bullets: list[str]) -> str:
        lines = [title.strip()] if title.strip() else []
        lines.extend(item.strip() for item in bullets if item.strip())
        return "\n".join(lines).strip()

    def _is_imported_project(self, project: Project) -> bool:
        project_metadata = project.project_metadata_json or {}
        return bool(project_metadata.get("is_imported"))

    def _recent_messages(self, project_id: str, limit: int = 12) -> list[dict[str, Any]]:
        stmt = (
            select(ProjectMessage)
            .where(ProjectMessage.project_id == project_id)
            .order_by(ProjectMessage.created_at.desc())
            .limit(limit)
        )
        messages = list(self.session.scalars(stmt))
        messages.reverse()
        return [
            {
                "role": item.role,
                "stage": item.stage,
                "scope_type": item.scope_type,
                "target_page_id": item.target_page_id,
                "content_md": item.content_md,
            }
            for item in messages
        ]

    def _build_system_decision(
        self,
        *,
        scope_type: str,
        target_stage: str,
        target_page_id: str | None,
        action_type: str,
        reason: str,
        execution_plan: list[dict[str, str]],
    ) -> dict[str, Any]:
        return {
            "scope_type": scope_type,
            "target_stage": target_stage,
            "target_page_id": target_page_id,
            "intent_type": action_type,
            "action_type": action_type,
            "should_execute": True,
            "needs_clarification": False,
            "requires_confirmation": False,
            "missing_data": [],
            "data_updates": {
                "question_patch": None,
                "answer_patch": None,
                "outline_patch": None,
                "page_patch": None,
                "summary_patch": None,
            },
            "execution_plan": execution_plan,
            "next_recommendations": [],
            "reason": reason,
        }

    def _page_action_execution_plan(self, action_type: str) -> list[dict[str, str]]:
        plans = {
            "page_generate_search_queries": [
                {"step_code": "page_generate_search_queries", "step_name": "生成页面搜索词", "reason": "把当前页结构化需求翻译成搜索词集合。"},
            ],
            "page_search_run": [
                {"step_code": "page_search_bocha", "step_name": "执行 Bocha 搜索", "reason": "先获取搜索摘要结果，再决定后续抓取。"},
                {"step_code": "page_search_read", "step_name": "抓取全文并写入资料池", "reason": "把搜索结果扩展成可引用正文。"},
                {"step_code": "page_search_vectorize", "step_name": "向量化资料池", "reason": "把正文切块并写入向量，供后续摘要和召回使用。"},
            ],
            "page_search_refresh": [
                {"step_code": "page_search_bocha", "step_name": "执行 Bocha 搜索", "reason": "先获取搜索摘要结果，再决定后续抓取。"},
                {"step_code": "page_search_read", "step_name": "抓取全文并写入资料池", "reason": "把搜索结果扩展成可引用正文。"},
                {"step_code": "page_search_vectorize", "step_name": "向量化资料池", "reason": "把正文切块并写入向量，供后续摘要和召回使用。"},
            ],
            "page_summary_generate": [
                {"step_code": "page_summary_generate", "step_name": "生成页面 summary", "reason": "只从当前页资料池内召回并生成摘要。"},
            ],
            "page_draft_generate": [
                {"step_code": "page_draft_generate", "step_name": "生成页面初稿", "reason": "基于当前页 summary 生成 draft。"},
            ],
            "page_design_generate": [
                {"step_code": "page_design_generate", "step_name": "生成页面设计稿", "reason": "基于当前页 draft 生成 design。"},
            ],
        }
        return plans.get(
            action_type,
            [{"step_code": action_type, "step_name": action_type, "reason": "按按钮动作执行。"}],
        )

    def _persist_agent_message(
        self,
        *,
        run: AgentRunRecorder,
        content_md: str,
        result_snapshot: dict[str, Any] | None = None,
    ) -> ProjectMessage:
        message = self._add_message(
            project_id=run.project.id,
            stage=run.stage,
            scope_type=run.scope_type,
            target_page_id=run.target_page_id,
            role="assistant",
            content_md=content_md,
            structured_payload_json={
                "message_kind": "agent_run",
                "agent_run_id": run.agent_run_id,
                "title": run.title,
                "router_decision": run.router_decision,
                "execution_plan": run.router_decision.get("execution_plan", []) if run.router_decision else [],
                "step_results": run.step_results,
                "next_recommendations": run.next_recommendations,
                "result_snapshot": result_snapshot or {},
            },
        )
        run.emit_message(message.id)
        return message

    def _format_exception_message(self, exc: Exception) -> str:
        detail = str(exc).strip()
        if detail:
            return f"{exc.__class__.__name__}: {detail}"
        return exc.__class__.__name__

    def _finalize_run_failure(
        self,
        *,
        run: AgentRunRecorder,
        step_code: str,
        step_name: str,
        exc: Exception,
        content_md: str,
        result_snapshot: dict[str, Any] | None = None,
    ) -> None:
        error_message = self._format_exception_message(exc)
        run.step_failed(step_code, step_name, error_message)
        snapshot = {
            "error_message": error_message,
            "error_type": exc.__class__.__name__,
            **(result_snapshot or {}),
        }
        self._persist_agent_message(
            run=run,
            content_md=content_md,
            result_snapshot=snapshot,
        )
        run.complete("failed")

    def _mark_page_stage_failed(self, page: ProjectPage, action_type: str) -> None:
        if action_type in {"page_search_run", "page_search_refresh"}:
            page.search_status = "failed"
        elif action_type == "page_summary_generate":
            page.summary_status = "failed"
        elif action_type == "page_draft_generate":
            page.draft_status = "failed"
        elif action_type == "page_design_generate":
            page.design_status = "failed"

    def run_bootstrap_flow(self, project_id: str) -> None:
        project = self._require_project(project_id)
        requirement_form = self._require_requirement_form(project)
        run = AgentRunRecorder(
            service=self,
            project=project,
            stage="init",
            scope_type="project",
            target_page_id=None,
            title="初始化资料准备",
            origin="system",
        )
        run.start()
        run.set_router_decision(
            self._build_system_decision(
                scope_type="project",
                target_stage="init",
                target_page_id=None,
                action_type="init_refresh_search",
                reason="项目创建后自动执行初始化搜索、建库和问题生成。",
                execution_plan=[
                    {"step_code": "I2", "step_name": "生成初始化搜索词", "reason": "把原始需求翻译成项目级查询。"},
                    {"step_code": "I3", "step_name": "执行 Bocha 搜索", "reason": "获取首轮搜索摘要。"},
                    {"step_code": "I4", "step_name": "生成页数推荐和问题", "reason": "基于搜索摘要快速产出固定项与问题。"},
                    {"step_code": "I5", "step_name": "抓取全文并写入 init_corpus", "reason": "把摘要结果扩展成可引用正文。"},
                    {"step_code": "I6", "step_name": "向量化 init_corpus", "reason": "把正文切块并建立后续召回能力。"},
                ],
            )
        )
        current_step_code = "I2"
        current_step_name = "生成初始化搜索词"
        try:
            run.step_started("I2", "生成初始化搜索词", "把原始需求翻译成项目级查询。")
            query_plan = self.research.build_query_plan(
                scope_type="project",
                session_role="init_discovery",
                request_text=project.request_text,
                project_stage="init",
                project_title=project.title,
                fixed_fields=self._build_fixed_field_values(project, requirement_form),
                answers=requirement_form.answers_json or {},
                latest_instruction=requirement_form.latest_instruction or "",
            )
            requirement_form.init_search_queries_json = query_plan
            run.data_updated(
                {
                    "entity": "requirement_form",
                    "update_kind": "init_queries",
                    "query_count": len(query_plan),
                }
            )
            run.step_completed("I2", "生成初始化搜索词", {"query_count": len(query_plan)})

            current_step_code = "I3"
            current_step_name = "执行 Bocha 搜索"
            run.step_started("I3", "执行 Bocha 搜索", "获取首轮搜索摘要。")

            def on_init_query_completed(payload: dict[str, Any]) -> None:
                requirement_form.init_search_results_json = payload["items"]
                run.data_updated(
                    {
                        "entity": "requirement_form",
                        "update_kind": "init_search_results",
                        "query_index": payload["query_index"],
                        "query_total": payload["query_total"],
                        "result_count": payload["result_count"],
                    }
                )
                run.step_progress(
                    "I3",
                    "执行 Bocha 搜索",
                    progress={
                        "current": payload["query_index"],
                        "total": payload["query_total"],
                        "label": f"已完成 {payload['query_index']}/{payload['query_total']} 条查询",
                    },
                    result={"result_count": payload["result_count"]},
                )

            search_results = self.research.search_query_summaries(
                query_plan,
                limit_per_query=4,
                on_query_completed=on_init_query_completed,
            )
            requirement_form.init_search_results_json = self.research.build_search_result_cards(search_results)
            run.step_completed("I3", "执行 Bocha 搜索", {"result_count": len(search_results)})

            current_step_code = "I4"
            current_step_name = "生成页数推荐和问题"
            run.step_started("I4", "生成页数推荐和问题", "基于搜索摘要快速产出固定项与问题。")
            package = self.generator.generate_init_fast_questions(
                project_title=project.title,
                request_text=project.request_text,
                init_search_results=search_results,
            )
            requirement_form.page_count_options_json = package["page_count_options"]
            requirement_form.ai_questions_json = package["ai_questions"]
            requirement_form.status = "ready"
            requirement_form.suggested_actions_json = [
                {
                    "code": "fill_required_fields",
                    "label": "补全固定项和问题答案",
                    "reason": "先完成页数、风格和问题答案，再进入大纲。",
                },
                {
                    "code": "refresh_init_search",
                    "label": "补充约束后重跑项目级搜索",
                    "reason": "如果首轮资料跑偏，可以要求重新搜索。",
                },
            ]
            run.data_updated(
                {
                    "entity": "requirement_form",
                    "update_kind": "init_questions",
                    "question_count": len(requirement_form.ai_questions_json),
                }
            )
            run.step_completed("I4", "生成页数推荐和问题", {"question_count": len(requirement_form.ai_questions_json)})

            current_step_code = "I5"
            current_step_name = "抓取全文并写入 init_corpus"
            run.step_started("I5", "抓取全文并写入 init_corpus", "把摘要结果扩展成可引用正文。")
            init_collection = self.research.get_or_create_init_collection(project)

            def on_init_read_progress(payload: dict[str, Any]) -> None:
                run.step_progress(
                    "I5",
                    "抓取全文并写入 init_corpus",
                    progress={
                        "current": payload["completed"],
                        "total": payload["total"],
                        "label": f"已处理 {payload['completed']}/{payload['total']} 条来源",
                    },
                    result={
                        "ingested_count": payload["ingested_count"],
                        "failed_count": payload["failed_count"],
                    },
                )

            candidate_sources, pending_chunk_records, read_summary = self.research.hydrate_search_results(
                collection=init_collection,
                search_results=search_results,
                replace=True,
                on_read_progress=on_init_read_progress,
            )
            candidate_sources = self.research.refresh_search_result_cards(candidate_sources)
            requirement_form.init_search_results_json = candidate_sources
            run.data_updated(
                {
                    "entity": "requirement_form",
                    "update_kind": "init_read",
                    "result_count": len(candidate_sources),
                    "read_ready": sum(1 for item in candidate_sources if item.get("read_status") in {"ready", "reused"}),
                    "read_failed": sum(1 for item in candidate_sources if item.get("read_status") == "failed"),
                }
            )
            run.step_completed(
                "I5",
                "抓取全文并写入 init_corpus",
                {
                    "result_count": len(candidate_sources),
                    "read_ready": sum(1 for item in candidate_sources if item.get("read_status") in {"ready", "reused"}),
                    "read_failed": sum(1 for item in candidate_sources if item.get("read_status") == "failed"),
                    **read_summary,
                },
            )

            current_step_code = "I6"
            current_step_name = "向量化 init_corpus"
            run.step_started("I6", "向量化 init_corpus", "把正文切块并建立后续召回能力。")

            def on_init_embedding_progress(payload: dict[str, Any]) -> None:
                run.step_progress(
                    "I6",
                    "向量化 init_corpus",
                    progress={
                        "current": payload["completed_chunks"],
                        "total": payload["total_chunks"],
                        "label": f"已写入 {payload['completed_chunks']}/{payload['total_chunks']} 个 chunk",
                    },
                    result={"document_count": payload["document_count"]},
                )

            embedding_stats = self.research.store_chunk_embeddings(
                pending_chunk_records,
                on_embedding_progress=on_init_embedding_progress,
            )
            candidate_sources = self.research.refresh_search_result_cards(candidate_sources)
            digest = self.research.build_collection_digest(init_collection.id)
            requirement_form.init_search_results_json = candidate_sources
            requirement_form.init_corpus_digest_json = digest
            run.data_updated(
                {
                    "entity": "requirement_form",
                    "update_kind": "init_vectorized",
                    "document_count": digest.get("document_count", 0),
                    "chunk_count": digest.get("chunk_count", 0),
                }
            )
            run.step_completed("I6", "向量化 init_corpus", {**digest, **embedding_stats})
            run.set_recommendations(requirement_form.suggested_actions_json)
            self._persist_agent_message(
                run=run,
                content_md="初始化资料已准备完成。现在可以填写页数、风格和补充问题答案；如果资料方向不对，也可以直接要求重跑项目级搜索。",
                result_snapshot={"requirement_form_status": requirement_form.status},
            )
            run.complete()
        except Exception as exc:
            requirement_form.status = "failed"
            self._finalize_run_failure(
                run=run,
                step_code=current_step_code,
                step_name=current_step_name,
                exc=exc,
                content_md="初始化资料准备失败。错误已经保留在当前动作卡片中，请先处理该错误再继续。",
                result_snapshot={"requirement_form_status": requirement_form.status},
            )

    def run_outline_flow(self, project_id: str) -> None:
        project = self._require_project(project_id)
        requirement_form = self._require_requirement_form(project)
        self._validate_requirement_form(project, requirement_form)
        run = AgentRunRecorder(
            service=self,
            project=project,
            stage="outline",
            scope_type="project",
            target_page_id=None,
            title="生成大纲并切换到搜索工作台",
            origin="system",
        )
        run.start()
        run.set_router_decision(
            self._build_system_decision(
                scope_type="project",
                target_stage="outline",
                target_page_id=None,
                action_type="outline_generate",
                reason="固定项齐备后生成大纲，并直接进入搜索工作台。",
                execution_plan=[
                    {"step_code": "O2", "step_name": "从 init_corpus 检索证据", "reason": "大纲只能使用项目级资料池。"},
                    {"step_code": "O3", "step_name": "生成大纲", "reason": "根据需求、固定项和证据生成章节与页面。"},
                    {"step_code": "O4", "step_name": "落库页面实体", "reason": "创建页面和首个版本。"},
                    {"step_code": "O5", "step_name": "切换到搜索页", "reason": "完成后直接进入搜索工作台，不自动搜索。"},
                ],
            )
        )
        fixed_fields = self._build_fixed_field_values(project, requirement_form)
        page_count_target = self._coerce_page_count(fixed_fields.get("page_count_target"), project.page_count_target) or 10
        project.page_count_target = page_count_target
        project.style_preset = str(fixed_fields.get("style_preset") or project.style_preset or "")
        current_step_code = "O2"
        current_step_name = "从 init_corpus 检索证据"
        try:
            run.step_started("O2", "从 init_corpus 检索证据", "大纲只能使用项目级资料池。")
            init_collection = self.research.get_or_create_init_collection(project)
            evidence_query_plan = self.research.build_query_plan(
                scope_type="project",
                session_role="outline_generate",
                request_text=project.request_text,
                project_stage="outline",
                project_title=project.title,
                fixed_fields=fixed_fields,
                answers=requirement_form.answers_json or {},
                latest_instruction=requirement_form.latest_instruction or "",
            )
            evidence_session = self.research.create_session(
                project_id=project.id,
                page_id=None,
                scope_type="project",
                session_role="outline_generate",
                research_goal="为大纲生成筛选项目级证据。",
                query_plan=evidence_query_plan,
                context_snapshot={"request_text": project.request_text, "fixed_fields": fixed_fields},
            )
            evidence = self.research.retrieve_for_collection(
                project=project,
                collection=init_collection,
                research_session=evidence_session,
                query_plan=evidence_query_plan,
                limit=200,
            )
            evidence_session.status = "completed" if evidence else "failed"
            run.step_completed("O2", "从 init_corpus 检索证据", {"citation_count": len(evidence)})

            current_step_code = "O3"
            current_step_name = "生成大纲"
            run.step_started("O3", "生成大纲", "根据需求、固定项和证据生成章节与页面。")
            outline_payload = self.generator.generate_outline(
                project_title=project.title,
                request_text=project.request_text,
                page_count_target=page_count_target,
                style_preset=project.style_preset or "",
                background_asset_path=project.background_asset_path,
                answers=requirement_form.answers_json or {},
                init_corpus_evidence=evidence,
            )
            run.step_completed("O3", "生成大纲", {"part_count": len(outline_payload["ppt_outline"].get("parts", []))})

            current_step_code = "O4"
            current_step_name = "落库页面实体"
            run.step_started("O4", "落库页面实体", "创建页面和首个版本。")
            self._rebuild_pages_from_outline(project, outline_payload)
            outline = OutlineVersion(
                project_id=project.id,
                version_no=(self.session.scalar(select(func.count(OutlineVersion.id)).where(OutlineVersion.project_id == project.id)) or 0) + 1,
                status="ready",
                outline_json=outline_payload,
            )
            self.session.add(outline)
            project.current_stage = "search"
            run.data_updated(
                {
                    "entity": "project",
                    "update_kind": "outline",
                    "page_count": len(self.list_pages(project.id)),
                }
            )
            run.step_completed("O4", "落库页面实体", {"page_count": len(self.list_pages(project.id))})

            current_step_code = "O5"
            current_step_name = "切换到搜索页"
            run.step_started("O5", "切换到搜索页", "完成后直接进入搜索工作台，不自动搜索。")
            run.status_changed({"current_stage": "search"})
            run.step_completed("O5", "切换到搜索页", {"current_stage": "search"})
            run.set_recommendations(
                [
                    {
                        "code": "page_generate_search_queries",
                        "label": "先为当前页生成搜索词",
                        "reason": "进入搜索页后默认不自动搜索，先看当前页职责是否正确。",
                    },
                    {
                        "code": "project_batch_search",
                        "label": "需要时再批量搜索",
                        "reason": "只有用户明确要求批量执行时才跑全项目。",
                    },
                ]
            )
            self._persist_agent_message(
                run=run,
                content_md="大纲生成完成，系统已进入搜索工作台。当前没有自动搜索任何页面，你可以先修改当前页标题和要点，再决定是否生成搜索词或执行搜索。",
                result_snapshot={"current_stage": "search"},
            )
            run.complete()
        except Exception as exc:
            self._finalize_run_failure(
                run=run,
                step_code=current_step_code,
                step_name=current_step_name,
                exc=exc,
                content_md="大纲生成失败。错误已经保留在当前动作卡片中。",
                result_snapshot={"current_stage": project.current_stage},
            )

    def _rebuild_pages_from_outline(self, project: Project, outline_payload: dict[str, Any]) -> None:
        for page in self.session.scalars(select(ProjectPage).where(ProjectPage.project_id == project.id)):
            self.session.delete(page)
        self.session.flush()

        page_defs: list[tuple[str, str | None, str, list[str]]] = []
        ppt_outline = outline_payload["ppt_outline"]
        page_defs.append(("cover", None, ppt_outline["cover"]["title"], ppt_outline["cover"].get("content", [])))
        page_defs.append(("toc", None, ppt_outline["table_of_contents"]["title"], ppt_outline["table_of_contents"].get("content", [])))
        for section in ppt_outline.get("parts", []):
            for page in section.get("pages", []):
                page_defs.append(("content", section["part_title"], page["title"], page.get("content", [])))
        page_defs.append(("end", None, ppt_outline["end_page"]["title"], ppt_outline["end_page"].get("content", [])))

        for sort_order, (role, part_title, title, content) in enumerate(page_defs, start=1):
            page = ProjectPage(
                project_id=project.id,
                page_code=f"page-{sort_order:02d}",
                page_role=role,
                part_title=part_title,
                sort_order=sort_order,
                outline_status="ready",
                search_status="confirmed" if role != "content" else "empty",
                summary_status="confirmed" if role != "content" else "empty",
                draft_status="empty",
                design_status="empty",
                page_summary_md="；".join(content[:2]) or title if role != "content" else "",
                page_summary_citations_json=[],
                page_search_queries_json=[],
                page_search_results_json=[],
                page_corpus_digest_json={},
                artifact_staleness_json={},
            )
            self.session.add(page)
            self.session.flush()
            self._new_page_brief_version(
                page=page,
                title=title,
                content_outline=content,
                section_title=part_title,
            )
            self._update_artifact_staleness(page)

    def _run_page_query_generation(
        self,
        *,
        project: Project,
        page: ProjectPage,
        latest_instruction: str,
    ) -> list[dict[str, str]]:
        if page.page_role != "content":
            raise RuntimeError("固定页不需要生成搜索词")
        brief = self._get_current_brief(page)
        if not brief:
            raise RuntimeError("页面结构不存在")
        queries = self.generator.generate_page_search_queries(
            project_title=project.title,
            project_request=project.request_text,
            page_id=page.id,
            page_title=brief.title,
            page_bullets=brief.content_outline_json,
            page_section_title=page.part_title,
            outline_full_snapshot=self._build_outline_snapshot(project.id),
            latest_instruction=latest_instruction,
        )
        page.page_search_queries_json = queries
        return queries

    def _run_page_search(
        self,
        *,
        project: Project,
        page: ProjectPage,
        latest_instruction: str,
        replace_existing: bool,
        run: AgentRunRecorder | None = None,
    ) -> dict[str, Any]:
        if page.page_role != "content":
            raise RuntimeError("固定页不需要页级搜索")
        queries = page.page_search_queries_json
        if not queries:
            if run is not None:
                run.step_started("page_search_prepare_queries", "补齐页面搜索词", "当前页还没有搜索词，先自动补齐。")
            queries = self._run_page_query_generation(
                project=project,
                page=page,
                latest_instruction=latest_instruction,
            )
            if run is not None:
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "search_queries",
                        "query_count": len(queries),
                    }
                )
                run.step_completed("page_search_prepare_queries", "补齐页面搜索词", {"query_count": len(queries)})
        page.search_status = "running"
        if run is not None:
            run.data_updated(
                {
                    "entity": "page",
                    "page_id": page.id,
                    "update_kind": "search_started",
                    "query_count": len(queries),
                }
            )
            run.step_started("page_search_bocha", "执行 Bocha 搜索", "先获取搜索摘要结果，再决定后续抓取。")

        def on_query_completed(payload: dict[str, Any]) -> None:
            page.page_search_results_json = payload["items"]
            if run is not None:
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "search_results",
                        "query_index": payload["query_index"],
                        "query_total": payload["query_total"],
                        "result_count": payload["result_count"],
                    }
                )
                run.step_progress(
                    "page_search_bocha",
                    "执行 Bocha 搜索",
                    progress={
                        "current": payload["query_index"],
                        "total": payload["query_total"],
                        "label": f"已完成 {payload['query_index']}/{payload['query_total']} 条查询",
                    },
                    result={"result_count": payload["result_count"]},
                )

        search_results = self.research.search_query_summaries(
            queries,
            limit_per_query=4,
            on_query_completed=on_query_completed if run is not None else None,
        )
        if run is not None:
            run.step_completed("page_search_bocha", "执行 Bocha 搜索", {"result_count": len(search_results)})
        collection = self.research.get_or_create_page_collection(project, page)
        if run is not None:
            run.step_started("page_search_read", "抓取全文并写入资料池", "把搜索结果扩展成可引用正文。")

        def on_read_progress(payload: dict[str, Any]) -> None:
            if run is None:
                return
            run.step_progress(
                "page_search_read",
                "抓取全文并写入资料池",
                progress={
                    "current": payload["completed"],
                    "total": payload["total"],
                    "label": f"已处理 {payload['completed']}/{payload['total']} 条来源",
                },
                result={
                    "ingested_count": payload["ingested_count"],
                    "failed_count": payload["failed_count"],
                },
            )

        candidate_sources, pending_chunk_records, read_summary = self.research.hydrate_search_results(
            collection=collection,
            search_results=search_results,
            replace=replace_existing,
            on_read_progress=on_read_progress if run is not None else None,
        )
        candidate_sources = self.research.refresh_search_result_cards(candidate_sources)
        page.page_search_results_json = candidate_sources
        read_ready = sum(1 for item in candidate_sources if item.get("read_status") in {"ready", "reused"})
        read_failed = sum(1 for item in candidate_sources if item.get("read_status") == "failed")
        if run is not None:
            run.data_updated(
                {
                    "entity": "page",
                    "page_id": page.id,
                    "update_kind": "search_read",
                    "result_count": len(candidate_sources),
                    "read_ready": read_ready,
                    "read_failed": read_failed,
                }
            )
            run.step_completed(
                "page_search_read",
                "抓取全文并写入资料池",
                {
                    "result_count": len(candidate_sources),
                    "read_ready": read_ready,
                    "read_failed": read_failed,
                },
            )

        if run is not None:
            run.step_started("page_search_vectorize", "向量化资料池", "把正文切块并写入向量，供后续摘要和召回使用。")

        def on_embedding_progress(payload: dict[str, Any]) -> None:
            if run is None:
                return
            run.step_progress(
                "page_search_vectorize",
                "向量化资料池",
                progress={
                    "current": payload["completed_chunks"],
                    "total": payload["total_chunks"],
                    "label": f"已写入 {payload['completed_chunks']}/{payload['total_chunks']} 个 chunk",
                },
                result={"document_count": payload["document_count"]},
            )

        embedding_stats = self.research.store_chunk_embeddings(
            pending_chunk_records,
            on_embedding_progress=on_embedding_progress if run is not None else None,
        )
        candidate_sources = self.research.refresh_search_result_cards(candidate_sources)
        digest = self.research.build_collection_digest(collection.id)
        session = self.research.create_session(
            project_id=project.id,
            page_id=page.id,
            scope_type="page",
            session_role="page_search",
            research_goal=f"为页面《{self._get_current_brief(page).title if self._get_current_brief(page) else page.page_code}》建立独立资料池。",
            query_plan=queries,
            context_snapshot={"latest_instruction": latest_instruction},
        )
        session.candidate_sources_json = candidate_sources
        session.status = "completed" if digest.get("document_count") else "failed"
        page.current_research_session_id = session.id
        page.page_search_results_json = candidate_sources
        page.page_corpus_digest_json = digest
        page.search_status = "ready" if digest.get("document_count") else "failed"
        page.summary_status = "stale" if page.page_summary_md else "empty"
        page.draft_status = "stale" if page.current_draft_version_id else "empty"
        page.design_status = "stale" if page.current_design_version_id else "empty"
        self._update_artifact_staleness(page)
        if run is not None:
            run.data_updated(
                {
                    "entity": "page",
                    "page_id": page.id,
                    "update_kind": "search_vectorized",
                    "document_count": digest.get("document_count", 0),
                    "chunk_count": digest.get("chunk_count", 0),
                }
            )
            run.step_completed(
                "page_search_vectorize",
                "向量化资料池",
                {
                    "document_count": digest.get("document_count", 0),
                    "chunk_count": digest.get("chunk_count", 0),
                    **embedding_stats,
                },
            )
        return {"query_count": len(queries), "result_count": len(candidate_sources), **digest}

    def _run_page_summary(
        self,
        *,
        project: Project,
        page: ProjectPage,
        latest_instruction: str,
    ) -> dict[str, Any]:
        if page.page_role != "content":
            raise RuntimeError("固定页不需要页级 summary 生成")
        brief = self._get_current_brief(page)
        if not brief:
            raise RuntimeError("页面结构不存在")
        if not page.page_corpus_digest_json.get("document_count"):
            if not self._is_imported_project(project):
                raise RuntimeError("当前页资料池为空，不能生成 summary")
            fallback_summary = self._build_imported_slide_summary(
                brief.title,
                brief.content_outline_json,
                page.page_role,
            )
            page.page_summary_md = fallback_summary
            page.page_summary_citations_json = []
            page.summary_status = "ready" if fallback_summary else "failed"
            page.draft_status = "stale" if page.current_draft_version_id else "empty"
            page.design_status = "stale" if page.current_design_version_id else "empty"
            self._update_artifact_staleness(page)
            return {"citation_count": 0, "summary_length": len(page.page_summary_md), "summary_source": "imported_outline"}
        page.summary_status = "running"
        query_plan = page.page_search_queries_json or [
            {
                "query_text": f"{brief.title} {' '.join(brief.content_outline_json)}".strip(),
                "query_purpose": "当前页核心事实和证据",
            }
        ]
        collection = self.research.get_or_create_page_collection(project, page)
        session = self.research.create_session(
            project_id=project.id,
            page_id=page.id,
            scope_type="page",
            session_role="page_summary",
            research_goal=f"从当前页资料池生成页面《{brief.title}》的详实摘要。",
            query_plan=query_plan,
            context_snapshot={"latest_instruction": latest_instruction},
        )
        selected = self.research.retrieve_for_collection(
            project=project,
            collection=collection,
            research_session=session,
            query_plan=query_plan,
            limit=20,
        )
        summary_package = self.generator.summarize_selected_sources(
            scope_type="page",
            research_goal=session.research_goal or "",
            selected_sources=selected,
        )
        session.summary_md = summary_package["summary_md"]
        session.status = "completed" if selected else "failed"
        page.current_research_session_id = session.id
        page.page_summary_md = summary_package["summary_md"]
        page.page_summary_citations_json = selected
        page.summary_status = "ready" if page.page_summary_md else "failed"
        page.draft_status = "stale" if page.current_draft_version_id else "empty"
        page.design_status = "stale" if page.current_design_version_id else "empty"
        self._update_artifact_staleness(page)
        return {"citation_count": len(selected), "summary_length": len(page.page_summary_md)}

    def _run_page_draft(
        self,
        *,
        project: Project,
        page: ProjectPage,
        latest_instruction: str,
    ) -> dict[str, Any]:
        brief = self._get_current_brief(page)
        if not brief:
            raise RuntimeError("页面结构不存在")
        summary_md = page.page_summary_md.strip()
        summary_source = "page_summary"
        if not summary_md and page.page_role != "content":
            summary_md = self._build_fixed_page_summary(brief.title, brief.content_outline_json)
            summary_source = "outline_brief"
        if not summary_md:
            raise RuntimeError("当前页 summary 为空，不能生成初稿")
        self._set_project_stage_at_least(project, "draft")
        page_context = {
            "page": {
                "page_id": page.id,
                "page_code": page.page_code,
                "page_role": page.page_role,
                "page_brief_version_id": brief.id,
                "title": brief.title,
                "content_outline": brief.content_outline_json,
                "content_summary": brief.content_summary,
            },
            "summary": {
                "summary_md": summary_md,
                "selected_sources": page.page_summary_citations_json,
            },
            "latest_instruction": latest_instruction,
        }
        svg = self.generator.generate_draft_svg(page_context=page_context)
        version_no = (self.session.scalar(select(func.count(DraftVersion.id)).where(DraftVersion.page_id == page.id)) or 0) + 1
        draft = DraftVersion(
            project_id=project.id,
            page_id=page.id,
            version_no=version_no,
            status="ready",
            page_brief_version_id=brief.id,
            research_session_id=page.current_research_session_id,
            draft_svg_markup=svg,
        )
        self.session.add(draft)
        self.session.flush()
        page.current_draft_version_id = draft.id
        page.draft_status = "ready"
        page.design_status = "stale" if page.current_design_version_id else "empty"
        self._update_artifact_staleness(page)
        return {"draft_version_id": draft.id, "summary_source": summary_source}

    def _run_page_design(
        self,
        *,
        project: Project,
        page: ProjectPage,
    ) -> dict[str, Any]:
        draft = self._get_current_draft(page)
        if not draft:
            raise RuntimeError("当前页初稿为空，不能生成设计稿")
        if not project.style_preset:
            raise RuntimeError("style_preset 未设置，不能生成设计稿")
        self._set_project_stage_at_least(project, "design")
        svg = self.generator.generate_design_svg(
            draft_svg=draft.draft_svg_markup,
            style_pack_id=project.style_preset,
            background_asset_path=project.background_asset_path,
        )
        version_no = (self.session.scalar(select(func.count(DesignVersion.id)).where(DesignVersion.page_id == page.id)) or 0) + 1
        design = DesignVersion(
            project_id=project.id,
            page_id=page.id,
            version_no=version_no,
            status="ready",
            draft_version_id=draft.id,
            style_pack_id=project.style_preset,
            background_asset_path=project.background_asset_path,
            design_svg_markup=svg,
        )
        self.session.add(design)
        self.session.flush()
        page.current_design_version_id = design.id
        page.design_status = "ready"
        self._update_artifact_staleness(page)
        return {"design_version_id": design.id}

    def run_page_action_flow(
        self,
        project_id: str,
        page_id: str,
        action_type: str,
        agent_run_id: str,
        replace_existing: bool,
    ) -> None:
        project = self._require_project(project_id)
        page = self._require_page(project_id, page_id)
        latest_instruction = ""
        run = AgentRunRecorder(
            service=self,
            project=project,
            stage=project.current_stage if project.current_stage != "outline" else "search",
            scope_type="page",
            target_page_id=page.id,
            title=f"执行页面动作：{action_type}",
            origin="button",
            message_id=None,
            agent_run_id=agent_run_id,
        )
        run.start()
        run.set_router_decision(
            self._build_system_decision(
                scope_type="page",
                target_stage=run.stage,
                target_page_id=page.id,
                action_type=action_type,
                reason="用户通过按钮直接触发页面动作。",
                execution_plan=self._page_action_execution_plan(action_type),
            )
        )
        result_snapshot: dict[str, Any] = {"page_id": page.id}
        step_name = action_type
        try:
            if action_type == "page_generate_search_queries":
                step_name = "生成页面搜索词"
                run.step_started(action_type, step_name, "把当前页结构化需求翻译成搜索词集合。")
                result_snapshot = {"queries": self._run_page_query_generation(project=project, page=page, latest_instruction=latest_instruction)}
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "search_queries",
                        "query_count": len(page.page_search_queries_json),
                    }
                )
                run.step_completed(action_type, step_name, {"query_count": len(page.page_search_queries_json)})
                message = f"已为当前页生成 {len(page.page_search_queries_json)} 条搜索词。下一步可以直接执行正式搜索。"
            elif action_type in {"page_search_run", "page_search_refresh"}:
                step_name = "执行页面搜索"
                result_snapshot = self._run_page_search(
                    project=project,
                    page=page,
                    latest_instruction=latest_instruction,
                    replace_existing=replace_existing or action_type == "page_search_refresh",
                    run=run,
                )
                message = "当前页资料池已更新。系统没有自动生成 summary，你可以继续手动生成 summary。"
            elif action_type == "page_summary_generate":
                step_name = "生成页面 summary"
                run.step_started(action_type, step_name, "只从当前页资料池内召回并生成摘要。")
                result_snapshot = self._run_page_summary(project=project, page=page, latest_instruction=latest_instruction)
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "summary",
                        "summary_length": result_snapshot.get("summary_length", 0),
                    }
                )
                run.step_completed(action_type, step_name, result_snapshot)
                message = "当前页 summary 已生成。系统没有自动继续出初稿。"
            elif action_type == "page_draft_generate":
                step_name = "生成页面初稿"
                run.step_started(action_type, step_name, "基于当前页 summary 生成 draft。")
                result_snapshot = self._run_page_draft(project=project, page=page, latest_instruction=latest_instruction)
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "draft",
                        "draft_version_id": result_snapshot.get("draft_version_id"),
                    }
                )
                run.step_completed(action_type, step_name, result_snapshot)
                message = "当前页初稿已生成。"
            elif action_type == "page_design_generate":
                step_name = "生成页面设计稿"
                run.step_started(action_type, step_name, "基于当前页 draft 和 style_preset 生成 design。")
                result_snapshot = self._run_page_design(project=project, page=page)
                run.data_updated(
                    {
                        "entity": "page",
                        "page_id": page.id,
                        "update_kind": "design",
                        "design_version_id": result_snapshot.get("design_version_id"),
                    }
                )
                run.step_completed(action_type, step_name, result_snapshot)
                message = "当前页设计稿已生成。"
            else:
                raise RuntimeError(f"不支持的页面动作: {action_type}")
        except Exception as exc:
            self._mark_page_stage_failed(page, action_type)
            self._update_artifact_staleness(page)
            self.session.commit()
            self._finalize_run_failure(
                run=run,
                step_code=action_type,
                step_name=step_name,
                exc=exc,
                content_md=f"{step_name}失败。错误已经保留在当前动作卡片中。",
                result_snapshot={"page_id": page.id},
            )
            return

        run.set_recommendations(self._default_recommendations_for_action(action_type))
        self._persist_agent_message(run=run, content_md=message, result_snapshot=result_snapshot)
        run.complete()

    def run_batch_action_flow(self, project_id: str, action_type: str, agent_run_id: str) -> None:
        project = self._require_project(project_id)
        run = AgentRunRecorder(
            service=self,
            project=project,
            stage=project.current_stage if project.current_stage != "outline" else "search",
            scope_type="project",
            target_page_id=None,
            title=f"批量执行：{action_type}",
            origin="button",
            agent_run_id=agent_run_id,
        )
        run.start()
        run.set_router_decision(
            self._build_system_decision(
                scope_type="project",
                target_stage=run.stage,
                target_page_id=None,
                action_type=action_type,
                reason="用户通过按钮明确触发批量动作。",
                execution_plan=[{"step_code": action_type, "step_name": action_type, "reason": "按批量规则逐页执行。"}],
            )
        )
        pages = list(self.session.scalars(select(ProjectPage).where(ProjectPage.project_id == project.id).order_by(ProjectPage.sort_order.asc())))
        processed = 0
        skipped = 0
        try:
            run.step_started(action_type, action_type, "按批量规则逐页执行。")
            for index, page in enumerate(pages, start=1):
                if action_type == "project_batch_search":
                    if page.page_role != "content":
                        skipped += 1
                    else:
                        self._run_page_search(project=project, page=page, latest_instruction="", replace_existing=True)
                        processed += 1
                elif action_type == "project_batch_summary":
                    if page.page_role != "content" or not page.page_corpus_digest_json.get("document_count"):
                        skipped += 1
                    else:
                        self._run_page_summary(project=project, page=page, latest_instruction="")
                        processed += 1
                elif action_type == "project_batch_draft":
                    if page.page_role == "content" and not page.page_summary_md:
                        skipped += 1
                    else:
                        self._run_page_draft(project=project, page=page, latest_instruction="")
                        processed += 1
                elif action_type == "project_batch_design":
                    if not page.current_draft_version_id:
                        skipped += 1
                    else:
                        self._run_page_design(project=project, page=page)
                        processed += 1
                else:
                    raise RuntimeError(f"不支持的批量动作: {action_type}")
                run.step_progress(
                    action_type,
                    action_type,
                    progress={
                        "current": index,
                        "total": len(pages),
                        "label": f"已扫描 {index}/{len(pages)} 页",
                    },
                    result={"processed": processed, "skipped": skipped},
                )
            run.step_completed(action_type, action_type, {"processed": processed, "skipped": skipped})
            run.set_recommendations(self._default_recommendations_for_action(action_type))
            self._persist_agent_message(
                run=run,
                content_md=f"批量动作 `{action_type}` 已执行完成。处理 {processed} 页，跳过 {skipped} 页。",
                result_snapshot={"processed": processed, "skipped": skipped},
            )
            run.complete()
        except Exception as exc:
            self._finalize_run_failure(
                run=run,
                step_code=action_type,
                step_name=action_type,
                exc=exc,
                content_md=f"批量动作 `{action_type}` 执行失败。错误已经保留在当前动作卡片中。",
                result_snapshot={"processed": processed, "skipped": skipped},
            )

    def run_message_flow(self, message_id: str) -> None:
        message = self.session.get(ProjectMessage, message_id)
        if not message:
            return
        project = self._require_project(message.project_id)
        page = self._require_page(project.id, message.target_page_id) if message.target_page_id else None
        ui_surface = str(message.structured_payload_json.get("ui_surface") or project.current_stage)
        run = AgentRunRecorder(
            service=self,
            project=project,
            stage=ui_surface if ui_surface in PROJECT_STAGE_ORDER else project.current_stage,
            scope_type=message.scope_type,
            target_page_id=message.target_page_id,
            title="处理聊天动作",
            origin="message",
            message_id=message.id,
        )
        run.start()
        try:
            requirement_form = self._require_requirement_form(project)
            router_payload = {
                "project_id": project.id,
                "project_stage": project.current_stage,
                "ui_surface": ui_surface,
                "latest_user_message": message.content_md,
                "recent_messages": self._recent_messages(project.id),
                "project_request": project.request_text,
                "workflow_constraints": project.workflow_constraints_json.get("items", []),
                "fixed_fields": self._build_fixed_field_values(project, requirement_form),
                "project_level_status_summary": self._build_project_level_status_summary(project),
                "outline_state_snapshot": self._build_outline_snapshot(project.id),
                "page_context": self._build_page_context_for_router(page),
                "default_scope_type": message.scope_type,
            }
            decision = self.generator.route_workspace_intent(router_payload=router_payload)
            run.set_router_decision(decision)
            run.set_recommendations(decision.get("next_recommendations", []))

            if decision["needs_clarification"] or decision["action_type"] == "reject":
                self._persist_agent_message(
                    run=run,
                    content_md=self._build_rejection_message(decision),
                    result_snapshot={"missing_data": decision.get("missing_data", [])},
                )
                run.complete("rejected")
                return

            action_type = decision["action_type"]
            if action_type == "init_refresh_search":
                self.run_bootstrap_flow(project.id)
                run.complete()
                return

            if action_type == "init_confirm_to_outline":
                try:
                    self.confirm_requirements(project.id, note_md=message.content_md)
                    self._persist_agent_message(
                        run=run,
                        content_md="初始化信息已满足要求，系统开始生成大纲。大纲完成后会直接进入搜索工作台。",
                        result_snapshot={"current_stage": "outline"},
                    )
                    run.complete()
                except Exception as exc:
                    self._finalize_run_failure(
                        run=run,
                        step_code=action_type,
                        step_name="确认初始化并生成大纲",
                        exc=exc,
                        content_md="确认初始化失败。错误已经保留在当前动作卡片中。",
                        result_snapshot={"current_stage": project.current_stage},
                    )
                return

            if action_type in {"project_batch_search", "project_batch_summary", "project_batch_draft", "project_batch_design"}:
                self.run_batch_action_flow(project.id, action_type, run.agent_run_id)
                run.complete()
                return

            if action_type in {"init_add_question", "init_update_question", "init_delete_question", "init_update_answer"}:
                result_snapshot: dict[str, Any]
                try:
                    if action_type == "init_update_answer":
                        run.step_started(action_type, "更新初始化答案", "根据聊天消息更新结构化答案，不自动重搜。")
                        answer_patch = decision["data_updates"].get("answer_patch")
                        if not isinstance(answer_patch, dict):
                            raise RuntimeError("router 没有返回可执行的 answer_patch")
                        result_snapshot = self._apply_init_answer_patch(project, requirement_form, answer_patch)
                        run.data_updated(
                            {
                                "entity": "requirement_form",
                                "update_kind": "init_answers",
                            }
                        )
                        run.step_completed(action_type, "更新初始化答案", result_snapshot)
                        content_md = "初始化答案已更新。系统没有自动重跑搜索，你可以继续修改，或明确要求重跑项目级搜索。"
                    else:
                        run.step_started("init_retrieval", "检索 init_corpus 证据", "问题增改前先从 init_corpus 做检索。")
                        init_collection = self.research.get_or_create_init_collection(project)
                        retrieval_query_plan = self.research.build_query_plan(
                            scope_type="project",
                            session_role="init_question_refine",
                            request_text=project.request_text,
                            project_stage="init",
                            project_title=project.title,
                            fixed_fields=self._build_fixed_field_values(project, requirement_form),
                            answers=requirement_form.answers_json or {},
                            latest_instruction=message.content_md,
                        )
                        refine_session = self.research.create_session(
                            project_id=project.id,
                            page_id=None,
                            scope_type="project",
                            session_role="init_question_refine",
                            research_goal="为初始化问题增删改提供项目级证据。",
                            query_plan=retrieval_query_plan,
                            context_snapshot={"latest_instruction": message.content_md},
                        )
                        evidence = self.research.retrieve_for_collection(
                            project=project,
                            collection=init_collection,
                            research_session=refine_session,
                            query_plan=retrieval_query_plan,
                            limit=200,
                        )
                        refine_session.status = "completed" if evidence else "failed"
                        run.step_completed("init_retrieval", "检索 init_corpus 证据", {"citation_count": len(evidence)})
                        question_patch = decision["data_updates"].get("question_patch")
                        if not isinstance(question_patch, dict):
                            raise RuntimeError("router 没有返回可执行的 question_patch")
                        result_snapshot = self._apply_init_question_patch(requirement_form, question_patch)
                        run.data_updated(
                            {
                                "entity": "requirement_form",
                                "update_kind": "init_questions",
                            }
                        )
                        run.step_completed(action_type, "更新初始化问题", result_snapshot)
                        content_md = "初始化问题集合已更新。当前不会自动重跑搜索；如果你要基于新问题重新看资料，请明确要求重跑项目级搜索。"
                    self._persist_agent_message(run=run, content_md=content_md, result_snapshot=result_snapshot)
                    run.complete()
                except Exception as exc:
                    self._finalize_run_failure(
                        run=run,
                        step_code=action_type,
                        step_name="更新初始化需求",
                        exc=exc,
                        content_md="初始化需求更新失败。错误已经保留在当前动作卡片中。",
                    )
                return

            if page is None:
                self._persist_agent_message(
                    run=run,
                    content_md="当前动作需要明确页面上下文，但这条消息没有绑定目标页。",
                    result_snapshot={"missing_data": ["target_page_id"]},
                )
                run.complete("rejected")
                return

            result_snapshot: dict[str, Any] = {"page_id": page.id}
            content_md = ""
            step_name = action_type
            try:
                if action_type == "page_update_outline_in_search":
                    step_name = "更新页面结构"
                    run.step_started(action_type, step_name, "修改标题、要点和章节归属，并只标记下游 stale。")
                    page_patch = self.generator.generate_page_outline_patch(
                        latest_user_message=message.content_md,
                        page_id=page.id,
                        page_title=self._get_current_brief(page).title if self._get_current_brief(page) else "",
                        page_bullets=self._get_current_brief(page).content_outline_json if self._get_current_brief(page) else [],
                        page_section_title=page.part_title,
                        outline_full_snapshot=self._build_outline_snapshot(project.id),
                    )
                    decision["data_updates"]["page_patch"] = page_patch
                    self.patch_page_outline(project.id, page.id, page_patch)
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "outline",
                        }
                    )
                    run.step_completed(action_type, step_name, page_patch)
                    result_snapshot = page_patch
                    content_md = f"当前页结构已更新：{page_patch.get('change_summary') or '标题和要点已写回数据库'}。系统没有自动重搜，相关下游产物已标记为 stale。"
                elif action_type == "page_generate_search_queries":
                    step_name = "生成页面搜索词"
                    run.step_started(action_type, step_name, "只重算当前页搜索词集合。")
                    result_snapshot = {"queries": self._run_page_query_generation(project=project, page=page, latest_instruction=message.content_md)}
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "search_queries",
                            "query_count": len(page.page_search_queries_json),
                        }
                    )
                    run.step_completed(action_type, step_name, {"query_count": len(page.page_search_queries_json)})
                    content_md = f"已为当前页生成 {len(page.page_search_queries_json)} 条搜索词。"
                elif action_type in {"page_search_run", "page_search_refresh"}:
                    step_name = "执行页面搜索"
                    result_snapshot = self._run_page_search(
                        project=project,
                        page=page,
                        latest_instruction=message.content_md,
                        replace_existing=action_type == "page_search_refresh",
                        run=run,
                    )
                    content_md = "当前页资料池已更新。系统没有自动继续生成 summary。"
                elif action_type == "page_summary_generate":
                    step_name = "生成页面 summary"
                    run.step_started(action_type, step_name, "只从当前页资料池生成摘要。")
                    result_snapshot = self._run_page_summary(project=project, page=page, latest_instruction=message.content_md)
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "summary",
                            "summary_length": result_snapshot.get("summary_length", 0),
                        }
                    )
                    run.step_completed(action_type, step_name, result_snapshot)
                    content_md = "当前页 summary 已生成。"
                elif action_type == "page_summary_edit":
                    step_name = "编辑页面 summary"
                    run.step_started(action_type, step_name, "根据用户要求改写当前页 summary，并标记 draft/design stale。")
                    patch = self.generator.generate_summary_patch(
                        latest_user_message=message.content_md,
                        page_title=self._get_current_brief(page).title if self._get_current_brief(page) else "",
                        page_bullets=self._get_current_brief(page).content_outline_json if self._get_current_brief(page) else [],
                        current_summary_md=page.page_summary_md,
                    )
                    decision["data_updates"]["summary_patch"] = patch
                    self.patch_page_summary(project.id, page.id, patch["summary_md"])
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "summary",
                            "summary_length": len(patch["summary_md"]),
                        }
                    )
                    run.step_completed(action_type, step_name, {"summary_length": len(patch["summary_md"])})
                    result_snapshot = patch
                    content_md = "当前页 summary 已按你的要求改写，draft/design 已标记为 stale。"
                elif action_type == "page_draft_generate":
                    step_name = "生成页面初稿"
                    run.step_started(action_type, step_name, "基于当前页 summary 生成 draft。")
                    result_snapshot = self._run_page_draft(project=project, page=page, latest_instruction=message.content_md)
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "draft",
                            "draft_version_id": result_snapshot.get("draft_version_id"),
                        }
                    )
                    run.step_completed(action_type, step_name, result_snapshot)
                    content_md = "当前页初稿已生成。"
                elif action_type == "page_design_generate":
                    step_name = "生成页面设计稿"
                    run.step_started(action_type, step_name, "基于当前页 draft 生成 design。")
                    result_snapshot = self._run_page_design(project=project, page=page)
                    run.data_updated(
                        {
                            "entity": "page",
                            "page_id": page.id,
                            "update_kind": "design",
                            "design_version_id": result_snapshot.get("design_version_id"),
                        }
                    )
                    run.step_completed(action_type, step_name, result_snapshot)
                    content_md = "当前页设计稿已生成。"
                else:
                    self._persist_agent_message(
                        run=run,
                        content_md="这条消息已经被识别到动作类型，但当前后端还没有对应执行器。",
                        result_snapshot={"action_type": action_type},
                    )
                    run.complete("rejected")
                    return
            except Exception as exc:
                self._mark_page_stage_failed(page, action_type)
                self._update_artifact_staleness(page)
                self.session.commit()
                self._finalize_run_failure(
                    run=run,
                    step_code=action_type,
                    step_name=step_name,
                    exc=exc,
                    content_md=f"{step_name}失败。错误已经保留在当前动作卡片中。",
                    result_snapshot={"page_id": page.id},
                )
                return

            self._persist_agent_message(run=run, content_md=content_md, result_snapshot=result_snapshot)
            run.complete()
        except Exception as exc:
            self._finalize_run_failure(
                run=run,
                step_code="route_workspace_intent",
                step_name="判断用户意图",
                exc=exc,
                content_md="处理聊天动作失败。错误已经保留在当前动作卡片中。",
                result_snapshot={
                    "message_id": message.id,
                    "page_id": page.id if page else None,
                    "ui_surface": ui_surface,
                },
            )

    def _default_recommendations_for_action(self, action_type: str) -> list[dict[str, Any]]:
        mapping = {
            "page_generate_search_queries": [
                {"code": "page_search_run", "label": "执行当前页搜索", "reason": "搜索词已经有了，下一步可以建立页级资料池。"}
            ],
            "page_search_run": [
                {"code": "page_summary_generate", "label": "生成当前页 summary", "reason": "资料池已经建立，现在可以只在当前页内做召回摘要。"}
            ],
            "page_summary_generate": [
                {"code": "page_draft_generate", "label": "生成当前页初稿", "reason": "summary 已准备好，可以继续出稿。"}
            ],
            "page_draft_generate": [
                {"code": "page_design_generate", "label": "生成当前页设计稿", "reason": "draft 已准备好，可以继续做设计增强。"}
            ],
            "project_batch_search": [
                {"code": "project_batch_summary", "label": "批量生成 summary", "reason": "批量搜索完成后，可以继续批量做页级摘要。"}
            ],
            "project_batch_summary": [
                {"code": "project_batch_draft", "label": "批量生成 draft", "reason": "所有页 summary 准备好后再批量出初稿。"}
            ],
            "project_batch_draft": [
                {"code": "project_batch_design", "label": "批量生成 design", "reason": "draft 已到位后可以继续批量设计。"}
            ],
        }
        return mapping.get(action_type, [])

    def _build_rejection_message(self, decision: dict[str, Any]) -> str:
        missing_data = decision.get("missing_data") or []
        if missing_data:
            return f"这条消息现在不能执行，缺少关键信息：{', '.join(missing_data)}。"
        return decision.get("reason") or "这条消息当前不能执行。"

    def _apply_init_answer_patch(self, project: Project, requirement_form: RequirementForm, answer_patch: dict[str, Any]) -> dict[str, Any]:
        question_code = str(answer_patch.get("question_code") or "").strip()
        if not question_code:
            raise RuntimeError("answer_patch 缺少 question_code")
        value = answer_patch.get("value")
        answers = dict(requirement_form.answers_json or {})
        answers[question_code] = value
        requirement_form.answers_json = answers
        if question_code == "page_count_target":
            project.page_count_target = self._coerce_page_count(value, project.page_count_target)
        if question_code == "style_preset" and value:
            project.style_preset = str(value)
        return {"question_code": question_code, "value": value}

    def _apply_init_question_patch(self, requirement_form: RequirementForm, question_patch: dict[str, Any]) -> dict[str, Any]:
        mode = str(question_patch.get("mode") or "upsert").strip()
        questions = list(requirement_form.ai_questions_json or [])
        if mode == "delete":
            question_code = str(question_patch.get("question_code") or "").strip()
            if not question_code:
                raise RuntimeError("question_patch 缺少 question_code")
            requirement_form.ai_questions_json = [item for item in questions if item.get("question_code") != question_code]
            answers = dict(requirement_form.answers_json or {})
            answers.pop(question_code, None)
            requirement_form.answers_json = answers
            return {"mode": "delete", "question_code": question_code}
        question = question_patch.get("question")
        if not isinstance(question, dict):
            raise RuntimeError("question_patch 缺少 question")
        question_code = str(question.get("question_code") or "").strip()
        if not question_code:
            raise RuntimeError("question_patch.question 缺少 question_code")
        filtered = [item for item in questions if item.get("question_code") != question_code]
        filtered.append(question)
        requirement_form.ai_questions_json = filtered
        return {"mode": "upsert", "question_code": question_code}


def run_bootstrap_job(project_id: str) -> None:
    with session_scope() as session:
        service = PptAgentService(session)
        service.run_bootstrap_flow(project_id)


def run_outline_job(project_id: str) -> None:
    with session_scope() as session:
        service = PptAgentService(session)
        service.run_outline_flow(project_id)


def run_page_action_job(
    project_id: str,
    page_id: str,
    action_type: str,
    agent_run_id: str,
    replace_existing: bool,
) -> None:
    with session_scope() as session:
        service = PptAgentService(session)
        service.run_page_action_flow(project_id, page_id, action_type, agent_run_id, replace_existing)


def run_batch_action_job(project_id: str, action_type: str, agent_run_id: str) -> None:
    with session_scope() as session:
        service = PptAgentService(session)
        service.run_batch_action_flow(project_id, action_type, agent_run_id)


def run_message_job(message_id: str) -> None:
    with session_scope() as session:
        service = PptAgentService(session)
        service.run_message_flow(message_id)
